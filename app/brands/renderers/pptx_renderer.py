"""
PPTX renderer for a Brand.

Takes an existing PPTX (built by pptx_builder) and applies the brand:
  1. COVER: dark primary-color bg + serif headline + accent eyebrow + meta-block
  2. INNER PAGES: top eyebrow + 'NN / TITLE' section label + serif H1 with
     accent rule underneath + footer with page numbers
"""
from __future__ import annotations

import logging
import time
from copy import deepcopy
from pathlib import Path
from typing import Any, Optional

from app.brands.base import Brand
from app.storage import (
    new_file_id,
    output_path,
    public_url,
)

log = logging.getLogger("aice.brand.pptx")


def _resolve(file_id: str) -> Path:
    p = output_path(file_id)
    if p.exists():
        return p
    from app.storage import input_path
    return input_path(file_id)


def _hex_to_rgb(hex_color: str):
    from pptx.dml.color import RGBColor
    s = hex_color.strip().lstrip("#")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _set_run(run, *, font=None, size_pt=None, bold=None, italic=None, color_hex=None):
    from pptx.util import Pt
    if font:
        run.font.name = font
    if size_pt is not None:
        run.font.size = Pt(size_pt)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic
    if color_hex:
        run.font.color.rgb = _hex_to_rgb(color_hex)


def _left_align_all(text_frame):
    """Set every paragraph in a text frame to left-align AND fix the textbox
    body properties so LibreOffice/PowerPoint render the alignment correctly.

    The default python-pptx add_textbox emits `<a:bodyPr wrap="none"><a:spAutoFit/>`,
    which auto-shrinks the shape to text width and (in LibreOffice) appears to
    center the text within the original geometry. We strip the autofit and set
    wrap="square" so the shape keeps its declared size and text aligns left."""
    from pptx.enum.text import PP_ALIGN
    from pptx.oxml.ns import qn

    # 1. Fix bodyPr: disable autofit + enable text wrapping inside the shape
    body_pr = text_frame._txBody.find(qn("a:bodyPr"))
    if body_pr is not None:
        for tag in ("a:spAutoFit", "a:noAutofit", "a:normAutofit"):
            for el in body_pr.findall(qn(tag)):
                body_pr.remove(el)
        body_pr.set("wrap", "square")

    # 2. Force left alignment on every paragraph
    for para in text_frame.paragraphs:
        para.alignment = PP_ALIGN.LEFT
        pPr = para._p.get_or_add_pPr()
        pPr.set("algn", "l")


def _send_to_back(shape):
    """Move a shape behind everything else on the slide so it acts as background."""
    sp = shape._element
    parent = sp.getparent()
    parent.remove(sp)
    parent.insert(2, sp)  # behind all real shapes (after nvGrpSpPr+grpSpPr)


def _strip_existing_chrome(slide):
    """Remove any text boxes / shapes that look like our previous theme chrome
    (footer text, page number, accent bars). We identify them by their position
    near the slide edges, since we don't tag them otherwise."""
    # In v1 we just leave existing shapes alone; brand renderer ADDS chrome.
    # Re-applying a brand a second time would stack — caller's responsibility
    # to not do that. Could add a marker xml-attr in a later version.
    pass


def _apply_cover(slide, brand: Brand, brief_meta: dict, slide_width, slide_height):
    """Re-render the cover slide in brand style: dark bg + serif headline +
    accent eyebrow + 4-column meta block."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN

    # Read the existing title/subtitle text BEFORE we re-render
    existing_title = ""
    existing_subtitle = ""
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        existing_title = slide.shapes.title.text_frame.text.strip()
    for ph in slide.placeholders:
        if ph is slide.shapes.title:
            continue
        if ph.has_text_frame:
            t = ph.text_frame.text.strip()
            if t and t != existing_title:
                existing_subtitle = t
                break

    # Use brief_meta override if provided
    title = brief_meta.get("title") or existing_title or "Cím"
    subtitle = brief_meta.get("subtitle") or existing_subtitle or ""

    # Wipe out any existing shapes on the cover (we're rebuilding it)
    sp_tree = slide.shapes._spTree  # type: ignore[attr-defined]
    for shape in list(slide.shapes):
        sp = shape._element
        sp_tree.remove(sp)

    # 1. Dark navy background — full-slide rectangle
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, slide_width, slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = _hex_to_rgb(brand.colors.primary)
    bg.line.fill.background()
    _send_to_back(bg)

    # 2. Top-left tiny org eyebrow (small caps, very subtle)
    org_name = brand.config.get("organization_name", "")
    if org_name and brand.motifs.inner_top_eyebrow:
        eb = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.4), slide_width - Inches(1.2), Inches(0.3)
        )
        eb.text_frame.text = org_name
        for run in eb.text_frame.paragraphs[0].runs or [eb.text_frame.paragraphs[0].add_run()]:
            _set_run(run, font=brand.fonts.effective_eyebrow(), size_pt=8.5,
                     color_hex=brand.colors.text_subtle)
            run.font.bold = True
        _left_align_all(eb.text_frame)

    # 3. Eyebrow band ("ELEMZÉSI ANYAG · HNI-TPL-01") — orange caps
    eyebrow_template = brand.config.get("cover_eyebrow_template", "{document_kind_label}")
    document_kind_label = brief_meta.get("document_kind_label",
                                         brand.config.get("default_document_kind_label", "DOKUMENTUM"))
    brief_code = brief_meta.get("brief_code", "")
    try:
        eyebrow_text = eyebrow_template.format(
            document_kind_label=document_kind_label, brief_code=brief_code,
        ).strip(" ·")
    except (KeyError, IndexError):
        eyebrow_text = document_kind_label

    eb2 = slide.shapes.add_textbox(
        Inches(0.6), Inches(2.1), slide_width - Inches(1.2), Inches(0.4)
    )
    eb2.text_frame.text = eyebrow_text.upper()
    for run in eb2.text_frame.paragraphs[0].runs or [eb2.text_frame.paragraphs[0].add_run()]:
        _set_run(run, font=brand.fonts.effective_eyebrow(), size_pt=10.5,
                 color_hex=brand.colors.accent)
        run.font.bold = True
    _left_align_all(eb2.text_frame)

    # 4. Sub-eyebrow (smaller, subtler — "UNIVERZÁLIS ELEMZÉSI RENDSZER" style)
    sub_eyebrow = brief_meta.get("sub_eyebrow", "")
    if sub_eyebrow:
        sb = slide.shapes.add_textbox(
            Inches(0.6), Inches(2.5), slide_width - Inches(1.2), Inches(0.3)
        )
        sb.text_frame.text = sub_eyebrow.upper()
        for run in sb.text_frame.paragraphs[0].runs or [sb.text_frame.paragraphs[0].add_run()]:
            _set_run(run, font=brand.fonts.effective_eyebrow(), size_pt=8.5,
                     color_hex=brand.colors.text_subtle)
        _left_align_all(sb.text_frame)

    # 5. Big serif headline (Georgia)
    headline = slide.shapes.add_textbox(
        Inches(0.6), Inches(3.0), slide_width - Inches(1.2), Inches(2.0)
    )
    headline.text_frame.word_wrap = True
    headline.text_frame.text = title
    for run in headline.text_frame.paragraphs[0].runs or [headline.text_frame.paragraphs[0].add_run()]:
        _set_run(run,
                 font=brand.fonts.heading,
                 size_pt=44 if len(title) <= 50 else 36,
                 bold=True,
                 color_hex=brand.colors.body_bg)
    _left_align_all(headline.text_frame)

    # 6. Subtitle (sans, white, 2-line lead)
    if subtitle:
        sb = slide.shapes.add_textbox(
            Inches(0.6), Inches(5.1), slide_width - Inches(1.2), Inches(1.0)
        )
        sb.text_frame.word_wrap = True
        sb.text_frame.text = subtitle
        for run in sb.text_frame.paragraphs[0].runs or [sb.text_frame.paragraphs[0].add_run()]:
            _set_run(run, font=brand.fonts.body, size_pt=14, color_hex=brand.colors.body_bg)
        _left_align_all(sb.text_frame)

    # 7. Thin gold rule above meta-block
    rule_top = slide_height - Inches(1.4)
    rule = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                                  Inches(0.6), rule_top, slide_width - Inches(1.2),
                                  int(Inches(0.012)))
    rule.fill.solid()
    rule.fill.fore_color.rgb = _hex_to_rgb(brand.colors.accent)
    rule.line.fill.background()

    # 8. 4-column meta block (DÁTUM | KOCKÁZATI SZINT | HORIZONT | DOKUMENTUM)
    if brand.motifs.cover_meta_block:
        fields = brand.config.get("cover_meta_fields",
                                  ["DÁTUM", "KOCKÁZATI SZINT", "HORIZONT", "DOKUMENTUM"])
        meta_values = brief_meta.get("meta_values", {})  # {field_label: value}
        col_count = len(fields)
        usable_w = slide_width - Inches(1.2)
        col_w = int(usable_w / col_count)
        meta_top = slide_height - Inches(1.1)
        for i, field_label in enumerate(fields):
            x = Inches(0.6) + col_w * i
            # Label (small caps, subtle gray)
            lbl = slide.shapes.add_textbox(x, meta_top, col_w, Inches(0.25))
            lbl.text_frame.text = field_label.upper()
            for run in lbl.text_frame.paragraphs[0].runs or [lbl.text_frame.paragraphs[0].add_run()]:
                _set_run(run, font=brand.fonts.effective_eyebrow(), size_pt=8,
                         color_hex=brand.colors.text_subtle)
                run.font.bold = True
            _left_align_all(lbl.text_frame)
            # Value (white, bold)
            val_text = meta_values.get(field_label, f"[{field_label.upper()}]")
            val = slide.shapes.add_textbox(x, meta_top + Inches(0.28), col_w, Inches(0.35))
            val.text_frame.text = str(val_text)
            for run in val.text_frame.paragraphs[0].runs or [val.text_frame.paragraphs[0].add_run()]:
                _set_run(run, font=brand.fonts.body, size_pt=11, bold=True,
                         color_hex=brand.colors.body_bg)
            _left_align_all(val.text_frame)


def _apply_inner_chrome(slide, brand: Brand, idx: int, total: int,
                       slide_width, slide_height,
                       document_label: str, section_number: Optional[int] = None,
                       section_title: Optional[str] = None):
    """Apply inner-page chrome: top eyebrow + section label + heading underline rule + footer."""
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    # 0. Replace the slide's title placeholder with a fresh textbox at a known
    #    position. The placeholder inherits geometry from the slide layout,
    #    which on the python-pptx blank theme overlaps with our section label
    #    chrome. A fresh textbox ignores that inheritance.
    title_text = ""
    if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
        title_text = slide.shapes.title.text_frame.text.strip()
        title_elem = slide.shapes.title._element
        title_elem.getparent().remove(title_elem)

    if title_text:
        title_box = slide.shapes.add_textbox(
            Inches(0.6), Inches(1.05), slide_width - Inches(1.2), Inches(0.75),
        )
        title_box.text_frame.word_wrap = True
        title_box.text_frame.text = title_text
        # Auto-shrink long titles
        sz = 26 if len(title_text) <= 40 else 22 if len(title_text) <= 60 else 18
        for run in title_box.text_frame.paragraphs[0].runs or [title_box.text_frame.paragraphs[0].add_run()]:
            _set_run(run, font=brand.fonts.heading, size_pt=sz, bold=True,
                     color_hex=brand.colors.text_primary)
        _left_align_all(title_box.text_frame)

    # 1. Top eyebrow (org name in subtle caps)
    if brand.motifs.inner_top_eyebrow:
        org = brand.config.get("organization_name", "")
        if org:
            eb = slide.shapes.add_textbox(Inches(0.6), Inches(0.3), slide_width - Inches(1.2), Inches(0.25))
            eb.text_frame.text = org
            for run in eb.text_frame.paragraphs[0].runs or [eb.text_frame.paragraphs[0].add_run()]:
                _set_run(run, font=brand.fonts.effective_eyebrow(), size_pt=8.5,
                         color_hex=brand.colors.text_subtle)
                run.font.bold = True
            _left_align_all(eb.text_frame)
            # Subtle horizontal rule under eyebrow
            rl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
                Inches(0.6), Inches(0.62), slide_width - Inches(1.2), int(Inches(0.008)))
            rl.fill.solid()
            rl.fill.fore_color.rgb = _hex_to_rgb(brand.colors.text_subtle)
            rl.line.fill.background()

    # 2. Section label (orange caps "NN / TITLE") — at y=0.78", above the title
    if section_number is not None and section_title:
        try:
            label_text = brand.motifs.section_label_format.format(
                number=section_number, title=section_title,
            )
        except (KeyError, IndexError):
            label_text = section_title
        if brand.motifs.section_label_caps:
            label_text = label_text.upper()
        sl = slide.shapes.add_textbox(Inches(0.6), Inches(0.78), slide_width - Inches(1.2), Inches(0.3))
        sl.text_frame.text = label_text
        for run in sl.text_frame.paragraphs[0].runs or [sl.text_frame.paragraphs[0].add_run()]:
            _set_run(run, font=brand.fonts.effective_eyebrow(), size_pt=10,
                     color_hex=brand.colors.accent)
            run.font.bold = True
        _left_align_all(sl.text_frame)

    # 3. Heading underline accent rule (orange thin line just below the title at y=1.85")
    if brand.motifs.heading_underline:
        rl = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE,
            Inches(0.6), Inches(1.92), Inches(1.5), int(Inches(0.025)))
        rl.fill.solid()
        rl.fill.fore_color.rgb = _hex_to_rgb(brand.colors.accent)
        rl.line.fill.background()

    # 4. Footer (document label + page number)
    if brand.motifs.page_numbers or document_label:
        try:
            footer_text = brand.motifs.footer_text_pattern.format(
                document_label=document_label, page=idx + 1, total=total,
            )
        except (KeyError, IndexError):
            footer_text = f"{document_label} · {idx + 1}"
        ft = slide.shapes.add_textbox(
            Inches(0.6), slide_height - Inches(0.4),
            slide_width - Inches(1.2), Inches(0.25))
        ft.text_frame.text = footer_text
        for run in ft.text_frame.paragraphs[0].runs or [ft.text_frame.paragraphs[0].add_run()]:
            _set_run(run, font=brand.fonts.effective_eyebrow(), size_pt=8.5,
                     color_hex=brand.colors.text_subtle)
        _left_align_all(ft.text_frame)


def _restyle_titles_and_body(slide, brand: Brand, is_cover: bool):
    """Apply brand fonts/colors to existing title + body text in the slide."""
    if is_cover:
        return  # cover is fully rerendered by _apply_cover
    if slide.shapes.title and slide.shapes.title.has_text_frame:
        for para in slide.shapes.title.text_frame.paragraphs:
            for run in para.runs or [para.add_run()]:
                _set_run(run,
                         font=brand.fonts.heading,
                         color_hex=brand.colors.text_primary,
                         bold=True)
    for shape in slide.placeholders:
        if shape is slide.shapes.title or not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs or [para.add_run()]:
                _set_run(run, font=brand.fonts.body,
                         color_hex=brand.colors.text_primary)
    # Also restyle non-placeholder textboxes (our explicit body textbox)
    title_shape = slide.shapes.title
    for shape in slide.shapes:
        if shape is title_shape:
            continue
        if shape in list(slide.placeholders):
            continue
        if not shape.has_text_frame:
            continue
        for para in shape.text_frame.paragraphs:
            for run in para.runs or [para.add_run()]:
                if run.font.name is None:
                    _set_run(run, font=brand.fonts.body,
                             color_hex=brand.colors.text_primary)


async def apply_brand_to_pptx(
    file_id: str, brand: Brand, brief_meta: Optional[dict] = None,
) -> dict:
    """Apply a Brand to a PPTX. brief_meta carries cover-specific values:
        - title, subtitle (override existing cover text)
        - document_kind_label (e.g. 'POLICY BRIEF')
        - brief_code (e.g. 'HNI-PB-2026-04')
        - sub_eyebrow (e.g. 'Európa energiapiac · 2026 Q2')
        - meta_values: {label: value} for the 4-col meta block (DÁTUM, etc.)
        - section_titles: optional [str] of titles per slide (idx-aligned, including cover)
        - document_label: footer label (defaults to brief_code or 'BRIEF')
    """
    from pptx import Presentation

    if brief_meta is None:
        brief_meta = {}
    document_label = (brief_meta.get("document_label")
                      or brief_meta.get("brief_code")
                      or "BRIEF")
    section_titles: list[str] = brief_meta.get("section_titles", []) or []

    src = _resolve(file_id)
    t0 = time.time()
    prs = Presentation(str(src))
    sw = prs.slide_width
    sh = prs.slide_height
    total = len(prs.slides)

    for idx, slide in enumerate(prs.slides):
        is_cover = idx == 0
        if is_cover:
            _apply_cover(slide, brand, brief_meta, sw, sh)
        else:
            # Section number = idx (1-based, since cover is 0)
            section_number = idx
            section_title = (
                section_titles[idx] if idx < len(section_titles)
                else (slide.shapes.title.text_frame.text.strip()
                      if slide.shapes.title and slide.shapes.title.has_text_frame
                      else None)
            )
            _restyle_titles_and_body(slide, brand, is_cover=False)
            _apply_inner_chrome(
                slide, brand, idx, total, sw, sh,
                document_label=document_label,
                section_number=section_number,
                section_title=section_title,
            )

    fid = new_file_id("pptx")
    dst = output_path(fid)
    prs.save(str(dst))
    return {
        "file_id": fid,
        "url": public_url(fid),
        "brand": brand.name,
        "slides_branded": total,
        "ms_elapsed": int((time.time() - t0) * 1000),
    }
