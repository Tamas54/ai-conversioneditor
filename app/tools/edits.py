"""
Deterministic edit operations for DOCX, HTML, MD.

These are the primitives that `edit_with_instruction` (LLM) compiles down to.
They are also exposed directly as MCP tools for cheap deterministic edits.

Supported scopes:
    - DOCX (python-docx): full read/write of paragraphs, headings, tables
    - HTML (BeautifulSoup): structural manipulation
    - MD (regex/line-based): simple but effective

PDF edits ALWAYS go through round-trip: PDF → DOCX → edit → DOCX → PDF.
The convert tool handles the round-trip; edits operate on DOCX/HTML/MD only.
"""
from __future__ import annotations

import logging
import re
import shutil
import time
from pathlib import Path
from typing import Optional

from app.storage import (
    input_path,
    new_file_id,
    output_path,
    public_url,
    temp_path,
)
from app.tools.formats import detect_from_path, normalize_format

log = logging.getLogger("aice.edits")


def _load_path(file_id: str) -> Path:
    p = input_path(file_id)
    if not p.exists():
        # also check outputs (a previous tool's output can be edited)
        p = output_path(file_id)
    if not p.exists():
        raise FileNotFoundError(f"file_id not found: {file_id}")
    return p


def _new_output_for(src: Path) -> tuple[str, Path]:
    fmt = detect_from_path(src)
    fid = new_file_id(fmt)
    return fid, output_path(fid)


# ============================================================
# DOCX edit primitives (python-docx)
# ============================================================


def _docx_find_replace(src: Path, dst: Path, pattern: str, replacement: str, regex: bool) -> int:
    """Returns the number of paragraphs in which a replacement was made.
    Note: pattern matching is paragraph-bounded — a pattern spanning two
    paragraphs will not match. For cross-paragraph edits, use insert_after
    or delete_section + insert_after."""
    from docx import Document

    doc = Document(str(src))
    rx = re.compile(pattern) if regex else None
    count = 0

    def replace_in_runs(paragraph):
        nonlocal count
        if not paragraph.runs:
            return
        full = "".join(r.text for r in paragraph.runs)
        if rx:
            new, n = rx.subn(replacement, full)
        else:
            n = full.count(pattern)
            new = full.replace(pattern, replacement) if n else full
        if n == 0 or new == full:
            return
        count += n
        paragraph.runs[0].text = new
        for r in paragraph.runs[1:]:
            r.text = ""

    for p in doc.paragraphs:
        replace_in_runs(p)
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    replace_in_runs(p)

    doc.save(str(dst))
    return count


def _docx_delete_section(src: Path, dst: Path, heading: str, level: Optional[int]):
    from docx import Document

    doc = Document(str(src))
    body = doc.element.body
    elements = list(body)

    # Find heading paragraph
    deleting = False
    found_level = None
    to_remove = []
    for el in elements:
        if el.tag.endswith("}p"):
            from docx.text.paragraph import Paragraph

            p = Paragraph(el, doc)
            style = (p.style.name or "") if p.style else ""
            is_heading = style.startswith("Heading")
            heading_lvl = None
            if is_heading:
                m = re.search(r"Heading\s+(\d)", style)
                if m:
                    heading_lvl = int(m.group(1))
            if not deleting and is_heading and p.text.strip() == heading.strip():
                if level is None or heading_lvl == level:
                    deleting = True
                    found_level = heading_lvl
                    to_remove.append(el)
                    continue
            if deleting and is_heading and heading_lvl is not None and found_level is not None:
                if heading_lvl <= found_level:
                    deleting = False
                    continue
        if deleting:
            to_remove.append(el)

    for el in to_remove:
        el.getparent().remove(el)
    doc.save(str(dst))


def _docx_insert_after(src: Path, dst: Path, anchor: str, content: str):
    """Insert `content` as a new paragraph immediately after the first
    paragraph whose text equals `anchor` (trimmed). Anchor can be any
    paragraph — heading or body — to give the agent maximum flexibility.
    For multi-line content, each line becomes its own paragraph."""
    from docx import Document

    doc = Document(str(src))
    target = anchor.strip()
    for p in doc.paragraphs:
        if p.text.strip() != target:
            continue
        anchor_elem = p._element
        # Insert in reverse so order is preserved when each line is added
        # immediately after `anchor_elem`.
        lines = content.splitlines() or [content]
        for line in reversed(lines):
            new_p = doc.add_paragraph(line)._element
            anchor_elem.addnext(new_p)
        doc.save(str(dst))
        return
    raise ValueError(f"anchor paragraph not found: {anchor!r}")


def _docx_replace_heading(src: Path, dst: Path, old_title: str, new_title: str):
    from docx import Document

    doc = Document(str(src))
    for p in doc.paragraphs:
        style = (p.style.name or "") if p.style else ""
        if style.startswith("Heading") and p.text.strip() == old_title.strip():
            for r in p.runs:
                r.text = ""
            if p.runs:
                p.runs[0].text = new_title
            else:
                p.add_run(new_title)
            doc.save(str(dst))
            return
    raise ValueError(f"Heading not found: {old_title!r}")


# ============================================================
# PDF in-place text replacement (PyMuPDF / fitz)
# ============================================================
# Strategy: search_for(pattern) → bbox per hit → redact (whiteout) →
# insert replacement at the same baseline with the original font/size/color.
# Preserves layout. Falls back to "helv" only when font file isn't found.

import shutil as _shutil
import subprocess as _subprocess


_FONT_FILE_CACHE: dict[str, Optional[str]] = {}


def _resolve_font_file(font_name: str) -> Optional[str]:
    """Map a PDF font name to a TTF/OTF file on the system. Returns None
    if no plausible match. Uses fc-match as last resort."""
    if font_name in _FONT_FILE_CACHE:
        return _FONT_FILE_CACHE[font_name]

    # Strip "AAAAAA+" subset prefix that embedded fonts often carry
    name = font_name.split("+")[-1]

    # Hand-curated candidates for common DOCX/LO/PDF fonts
    base_dirs = [
        "/usr/share/fonts/truetype/liberation",
        "/usr/share/fonts/truetype/dejavu",
        "/usr/share/fonts/truetype/freefont",
        "/usr/share/fonts/truetype/noto",
    ]
    candidates: list[str] = []
    # Try with dash-stripped variants too
    name_variants = {name, name.replace("-", ""), name.replace(" ", "")}
    for d in base_dirs:
        for n in name_variants:
            candidates.extend([
                f"{d}/{n}.ttf", f"{d}/{n}.otf",
                f"{d}/{n}-Regular.ttf",
            ])
    for c in candidates:
        if Path(c).exists():
            _FONT_FILE_CACHE[font_name] = c
            return c

    # fc-match fallback — handles arbitrary font names robustly
    try:
        r = _subprocess.run(
            ["fc-match", "-f", "%{file}", name],
            capture_output=True, text=True, timeout=2,
        )
        if r.returncode == 0 and r.stdout.strip():
            _FONT_FILE_CACHE[font_name] = r.stdout.strip()
            return r.stdout.strip()
    except (FileNotFoundError, _subprocess.TimeoutExpired):
        pass

    _FONT_FILE_CACHE[font_name] = None
    return None


def _pdf_color_to_rgb(c: int | tuple) -> tuple[float, float, float]:
    if isinstance(c, tuple):
        return tuple(float(x) for x in c[:3])
    return (((c >> 16) & 0xFF) / 255.0, ((c >> 8) & 0xFF) / 255.0, (c & 0xFF) / 255.0)


def _parse_to_unicode_cmap(cmap_bytes: bytes) -> dict[int, str]:
    """Parse a ToUnicode CMap stream into {cid_int: unicode_str}. Handles
    bfchar (single mappings) and bfrange (ranges + explicit lists)."""
    text = cmap_bytes.decode("latin-1", errors="replace")
    cid_to_uni: dict[int, str] = {}
    for block in re.findall(r"beginbfchar(.*?)endbfchar", text, re.S):
        for m in re.finditer(r"<([0-9A-Fa-f]+)>\s*<([0-9A-Fa-f]+)>", block):
            try:
                cid = int(m.group(1), 16)
                uni = bytes.fromhex(m.group(2)).decode("utf-16-be")
                cid_to_uni[cid] = uni
            except (ValueError, UnicodeDecodeError):
                continue
    for block in re.findall(r"beginbfrange(.*?)endbfrange", text, re.S):
        # form 1: <s> <e> <u_start>
        for m in re.finditer(r"<([0-9A-Fa-f]+)>\s*<([0-9A-Fa-f]+)>\s*<([0-9A-Fa-f]+)>", block):
            try:
                s, e, u = int(m.group(1), 16), int(m.group(2), 16), int(m.group(3), 16)
                for i, cid in enumerate(range(s, e + 1)):
                    cid_to_uni[cid] = chr(u + i)
            except (ValueError, OverflowError):
                continue
        # form 2: <s> <e> [<u1> <u2> ...]
        for m in re.finditer(r"<([0-9A-Fa-f]+)>\s*<([0-9A-Fa-f]+)>\s*\[([^\]]*)\]", block):
            try:
                s, e = int(m.group(1), 16), int(m.group(2), 16)
                unis = re.findall(r"<([0-9A-Fa-f]+)>", m.group(3))
                for i, cid in enumerate(range(s, e + 1)):
                    if i < len(unis):
                        try:
                            cid_to_uni[cid] = bytes.fromhex(unis[i]).decode("utf-16-be")
                        except (ValueError, UnicodeDecodeError):
                            pass
            except (ValueError, OverflowError):
                continue
    return cid_to_uni


def _cmap_decode(b: bytes, cmap: dict[int, str], cid_size: int = 2) -> str:
    return "".join(cmap.get(int.from_bytes(b[i:i + cid_size], "big"), "") for i in range(0, len(b), cid_size))


def _cmap_encode(s: str, cmap: dict[int, str], cid_size: int = 2) -> Optional[bytes]:
    inv = {v: k for k, v in cmap.items()}
    out = bytearray()
    for ch in s:
        cid = inv.get(ch)
        if cid is None:
            return None
        out.extend(cid.to_bytes(cid_size, "big"))
    return bytes(out)


def _pdf_find_replace_content_stream(src: Path, dst: Path, pattern: str, replacement: str) -> int:
    """Bytestream-level text replacement: patches Tj/TJ operands directly,
    preserving the original embedded font, kerning, and positioning. Iterates
    every CMap on the page and tries each — uses whichever decodes the pattern
    successfully. Falls back to the redact+insert approach for hits this method
    can't patch (e.g. patterns split across multiple Tj operators)."""
    import fitz

    doc = fitz.open(str(src))
    total_patched = 0
    unpatched_rects_per_page: dict[int, list] = {}

    for page_idx, page in enumerate(doc):
        rects = page.search_for(pattern)
        if not rects:
            continue

        # Build CMaps for all fonts on the page
        cmaps: list[dict[int, str]] = []
        for xref, *_ in page.get_fonts():
            fdict = doc.xref_object(xref)
            m = re.search(r"/ToUnicode\s+(\d+)\s+0\s+R", fdict)
            if not m:
                continue
            cb = doc.xref_stream(int(m.group(1)))
            if cb:
                cm = _parse_to_unicode_cmap(cb)
                if cm:
                    cmaps.append(cm)

        if not cmaps:
            unpatched_rects_per_page[page_idx] = list(rects)
            continue

        contents_xrefs = page.get_contents()
        merged = b"".join(doc.xref_stream(x) for x in contents_xrefs)
        text = merged.decode("latin-1", errors="replace")

        op_re = re.compile(
            r"<(?P<tj>[0-9A-Fa-f]+)>\s*Tj"
            r"|\[(?P<tjarr>(?:<[0-9A-Fa-f]+>|-?\d+\.?\d*\s*)+)\]\s*TJ"
        )
        replacements: list[tuple[int, int, bytes]] = []
        page_patched = 0

        for m in op_re.finditer(text):
            if m.group("tj") is not None:
                cb = bytes.fromhex(m.group("tj"))
                # Try each cmap; pick the one whose decoded string contains the pattern
                for cmap in cmaps:
                    decoded = _cmap_decode(cb, cmap)
                    if pattern not in decoded:
                        continue
                    new_dec = decoded.replace(pattern, replacement)
                    nb = _cmap_encode(new_dec, cmap)
                    if nb is None:
                        continue
                    replacements.append(
                        (m.start(), m.end(), f"<{nb.hex().upper()}>Tj".encode("latin-1"))
                    )
                    page_patched += decoded.count(pattern)
                    break
            else:
                arr_text = m.group("tjarr")
                pieces: list[tuple[str, Any]] = []
                for tok in re.finditer(r"<([0-9A-Fa-f]+)>|(-?\d+\.?\d*)", arr_text):
                    if tok.group(1):
                        pieces.append(("s", bytes.fromhex(tok.group(1))))
                    else:
                        pieces.append(("n", tok.group(2)))
                # Try each cmap on the concatenated string content
                for cmap in cmaps:
                    str_pieces_idx = [i for i, (k, _) in enumerate(pieces) if k == "s"]
                    if not str_pieces_idx:
                        continue
                    str_decoded = [_cmap_decode(pieces[i][1], cmap) for i in str_pieces_idx]
                    full = "".join(str_decoded)
                    if pattern not in full:
                        continue
                    new_full = full.replace(pattern, replacement)
                    if new_full == full:
                        continue
                    # Determine which string-pieces span the matches.
                    pos_to_str_piece: list[tuple[int, int, int]] = []
                    cur = 0
                    for spi, dec in zip(str_pieces_idx, str_decoded):
                        pos_to_str_piece.append((cur, cur + len(dec), spi))
                        cur += len(dec)
                    pi_ranges: set[tuple[int, int]] = set()
                    start = 0
                    while True:
                        i = full.find(pattern, start)
                        if i == -1:
                            break
                        end = i + len(pattern)
                        sa = sb = None
                        for ps, pe, spi in pos_to_str_piece:
                            if ps <= i < pe and sa is None:
                                sa = spi
                            if ps < end <= pe:
                                sb = spi
                                break
                        if sa is None or sb is None:
                            sb = pos_to_str_piece[-1][2]
                            sa = sa or pos_to_str_piece[0][2]
                        pi_ranges.add((sa, sb))
                        start = end
                    pi_ranges_sorted = sorted(pi_ranges)
                    out_pieces: list[tuple[str, Any]] = []
                    i_p = 0
                    range_iter = iter(pi_ranges_sorted)
                    next_range = next(range_iter, None)
                    while i_p < len(pieces):
                        if next_range and pieces[i_p][0] == "s" and i_p == next_range[0]:
                            sa, sb = next_range
                            buf = ""
                            j = i_p
                            while j <= sb and j < len(pieces):
                                if pieces[j][0] == "s":
                                    buf += _cmap_decode(pieces[j][1], cmap)
                                # numbers inside the matched span are dropped (kerning between
                                # letters of the contiguous word being replaced)
                                j += 1
                            new_buf = buf.replace(pattern, replacement)
                            nb = _cmap_encode(new_buf, cmap)
                            if nb is None:
                                out_pieces.extend(pieces[i_p:sb + 1])
                            else:
                                out_pieces.append(("s", nb))
                            i_p = sb + 1
                            next_range = next(range_iter, None)
                        else:
                            out_pieces.append(pieces[i_p])
                            i_p += 1
                    parts = [
                        f"<{v.hex().upper()}>" if k == "s" else str(v)
                        for k, v in out_pieces
                    ]
                    new_op = ("[" + " ".join(parts) + "]TJ").encode("latin-1")
                    replacements.append((m.start(), m.end(), new_op))
                    page_patched += full.count(pattern)
                    break

        if replacements:
            new_stream = bytearray(merged)
            for s, e, new_op in sorted(replacements, key=lambda x: -x[0]):
                new_stream[s:e] = new_op
            doc.update_stream(contents_xrefs[0], bytes(new_stream))
            for x in contents_xrefs[1:]:
                doc.update_stream(x, b"")

        total_patched += page_patched
        # Hits that survived (the byte-stream patch couldn't reach them)
        # need a second pass via the redact+insert fallback
        if page_patched < len(rects):
            # Recompute rects on the (now-patched) page; remaining hits = still need fix
            still = [r for r in page.search_for(pattern)]
            if still:
                unpatched_rects_per_page[page_idx] = still

    # --- Fallback for remaining hits: redact + insert via PyMuPDF.
    #     Use the EMBEDDED font from the source PDF (via doc.extract_font) when
    #     possible — this produces pixel-perfect rendering since it's literally
    #     the same TTF bytes that the rest of the page uses.
    import fitz as _fitz
    fallback_used = 0

    # Map font_name -> (alias, fontbuffer) for embedded fonts extracted from doc
    embedded_font_cache: dict[str, tuple[str, bytes]] = {}

    def _get_embedded_font(font_name: str, page) -> Optional[tuple[str, bytes]]:
        """Try to find this font name among the page's fonts and extract its bytes."""
        if not font_name:
            return None
        if font_name in embedded_font_cache:
            return embedded_font_cache[font_name]
        # Strip subset prefix (e.g. "VRMYHF+Liberation-Sans" → match by suffix)
        base = font_name.split("+")[-1]
        for xref, ext, ftype, basefont, fname, encoding in page.get_fonts():
            xbase = (basefont or "").split("+")[-1]
            if xbase == base or basefont == font_name or fname == font_name:
                try:
                    fb = doc.extract_font(xref)
                    # extract_font returns (name, ext, type, buffer)
                    if fb and len(fb) >= 4 and fb[3]:
                        alias = f"pdf_emb_{xref}"
                        embedded_font_cache[font_name] = (alias, fb[3])
                        return embedded_font_cache[font_name]
                except Exception:
                    pass
        embedded_font_cache[font_name] = None  # type: ignore
        return None

    for page_idx, rects in unpatched_rects_per_page.items():
        page = doc[page_idx]
        text_dict = page.get_text("dict")
        spans = [
            (_fitz.Rect(sp["bbox"]), sp)
            for blk in text_dict.get("blocks", []) if blk.get("type") == 0
            for ln in blk.get("lines", [])
            for sp in ln.get("spans", [])
        ]
        plans = []
        for r in rects:
            match_span = next((sp for srect, sp in spans if srect.intersects(r)), None)
            plans.append((r, match_span))
        for r, _ in plans:
            page.add_redact_annot(r, fill=(1, 1, 1))
        try:
            page.apply_redactions()
        except Exception:
            continue
        # Pre-register fonts on this page so insert_text can reference them by alias
        registered_aliases_on_page: set[str] = set()
        for r, span in plans:
            font_name = (span or {}).get("font", "")
            emb = _get_embedded_font(font_name, page)
            if emb:
                alias, buf = emb
                if alias not in registered_aliases_on_page:
                    try:
                        page.insert_font(fontname=alias, fontbuffer=buf)
                        registered_aliases_on_page.add(alias)
                    except Exception as e:
                        log.debug("insert_font(%s) failed: %s", alias, e)

        for r, span in plans:
            size = (span or {}).get("size", 10)
            color = _pdf_color_to_rgb((span or {}).get("color", 0))
            font_name = (span or {}).get("font", "")
            kw: dict[str, Any] = {"fontsize": size, "color": color}
            emb = _get_embedded_font(font_name, page)
            if emb and emb[0] in registered_aliases_on_page:
                kw["fontname"] = emb[0]
            else:
                ff = _resolve_font_file(font_name) if font_name else None
                if ff:
                    alias = f"emb_{abs(hash(ff)) % (10 ** 8)}"
                    kw["fontname"] = alias
                    kw["fontfile"] = ff
            page.insert_text(
                (r.x0, r.y1 - max(1.0, size * 0.18)), replacement, **kw,
            )
            fallback_used += 1

    if fallback_used:
        log.info("PDF find_replace: %d hits patched in-stream, %d via fallback redact+insert",
                 total_patched, fallback_used)
        try:
            doc.subset_fonts(verbose=False)
        except Exception:
            pass

    doc.save(str(dst), deflate=True, garbage=4, clean=True)
    doc.close()
    return total_patched + fallback_used


# Backwards-compat alias (older calls used the redact-only function name)
_pdf_find_replace_inplace = _pdf_find_replace_content_stream


# ============================================================
# DOCX styling primitives (color, font, highlight, etc.)
# ============================================================


_HIGHLIGHT_NAME_TO_INDEX = {
    "yellow": 7, "green": 4, "cyan": 3, "magenta": 5, "blue": 2, "red": 6,
    "darkblue": 9, "darkcyan": 10, "darkgreen": 11, "darkmagenta": 12, "darkred": 13,
    "darkyellow": 14, "darkgray": 15, "lightgray": 16, "black": 1, "white": 8,
}


def _parse_rgb(color: str):
    """Accept '#RRGGBB' or 'RRGGBB' or named CSS-ish color → docx RGBColor."""
    from docx.shared import RGBColor

    named = {
        "black": "000000", "white": "FFFFFF", "red": "C0392B", "green": "27AE60",
        "blue": "2C3E50", "navy": "1F3A5F", "gold": "B8860B", "orange": "E67E22",
        "purple": "8E44AD", "teal": "16A085", "gray": "7F8C8D", "darkgray": "555555",
        "lightgray": "BDBDBD",
        # extras for the dashboard's gold theme
        "darkblue": "1F3A5F", "darkgreen": "1B5E20", "darkred": "8B1F1F",
    }
    s = color.strip().lstrip("#")
    if s.lower() in named:
        s = named[s.lower()]
    if len(s) != 6:
        raise ValueError(f"invalid color: {color!r}")
    return RGBColor(int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


def _docx_apply_style_to_paragraph(p, *, color=None, highlight=None, bold=None,
                                   italic=None, underline=None, font_size=None,
                                   font_family=None):
    from docx.shared import Pt
    from docx.enum.text import WD_COLOR_INDEX

    rgb = _parse_rgb(color) if color else None
    hl = None
    if highlight:
        idx = _HIGHLIGHT_NAME_TO_INDEX.get(highlight.lower())
        if idx is None:
            raise ValueError(f"unknown highlight color: {highlight}")
        hl = WD_COLOR_INDEX(idx)
    for r in p.runs or [p.add_run("")]:
        if rgb is not None:
            r.font.color.rgb = rgb
        if hl is not None:
            r.font.highlight_color = hl
        if bold is not None:
            r.bold = bold
        if italic is not None:
            r.italic = italic
        if underline is not None:
            r.underline = underline
        if font_size is not None:
            r.font.size = Pt(font_size)
        if font_family:
            r.font.name = font_family


def _docx_style_paragraph(src: Path, dst: Path, anchor: str, **style) -> int:
    """Apply formatting to every paragraph whose text equals `anchor`.
    Returns count of paragraphs styled."""
    from docx import Document

    doc = Document(str(src))
    target = anchor.strip()
    n = 0
    for p in doc.paragraphs:
        if p.text.strip() == target:
            _docx_apply_style_to_paragraph(p, **style)
            n += 1
    for table in doc.tables:
        for row in table.rows:
            for cell in row.cells:
                for p in cell.paragraphs:
                    if p.text.strip() == target:
                        _docx_apply_style_to_paragraph(p, **style)
                        n += 1
    if n == 0:
        raise ValueError(f"no paragraph found with text: {anchor!r}")
    doc.save(str(dst))
    return n


def _docx_style_all_headings(src: Path, dst: Path, *, levels=None, **style) -> int:
    """Apply formatting to all heading paragraphs (optionally filtered to specific levels)."""
    from docx import Document

    doc = Document(str(src))
    levels_set = set(levels) if levels else None
    n = 0
    for p in doc.paragraphs:
        sname = (p.style.name or "") if p.style else ""
        if not sname.startswith("Heading"):
            continue
        if levels_set is not None:
            m = re.search(r"Heading\s+(\d)", sname)
            if not m or int(m.group(1)) not in levels_set:
                continue
        _docx_apply_style_to_paragraph(p, **style)
        n += 1
    doc.save(str(dst))
    return n


def _docx_insert_image_after(src: Path, dst: Path, anchor: str, image_path: Path,
                             width_inches: Optional[float] = None) -> None:
    from docx import Document
    from docx.shared import Inches

    doc = Document(str(src))
    target = anchor.strip()
    for p in doc.paragraphs:
        if p.text.strip() != target:
            continue
        # Insert a new empty paragraph right after the anchor, then add the picture into it
        from copy import deepcopy
        new_para_xml = deepcopy(p._element)
        # Strip child runs from clone — we just want a fresh paragraph element
        for child in list(new_para_xml):
            if child.tag.endswith("}r") or child.tag.endswith("}pPr"):
                # keep pPr but blank runs
                if child.tag.endswith("}r"):
                    new_para_xml.remove(child)
        p._element.addnext(new_para_xml)
        # Wrap the inserted XML element back into a python-docx Paragraph
        from docx.text.paragraph import Paragraph
        new_p = Paragraph(new_para_xml, p._parent)
        run = new_p.add_run()
        kwargs = {"width": Inches(width_inches)} if width_inches else {}
        run.add_picture(str(image_path), **kwargs)
        doc.save(str(dst))
        return
    raise ValueError(f"anchor paragraph not found: {anchor!r}")


def _attach_svg_to_last_picture(docx_path: Path, svg_path: Path) -> None:
    """Post-process a DOCX to attach an SVG vector layer alongside the
    most recently inserted PNG picture.

    Modern Word (2016+) and LibreOffice 7+ render the SVG (vector, crisp at
    any zoom and in PDF export); older readers fall back to the PNG that
    python-docx already embedded. This is Microsoft's official extension
    mechanism — no DOCX schema violation.

    Patches three parts of the DOCX zip:
      1. [Content_Types].xml — register image/svg+xml default MIME
      2. word/_rels/document.xml.rels — new rId pointing at the SVG file
      3. word/document.xml — inject <a:extLst><a:ext uri="{96D…}"
         ><asvg:svgBlip r:embed=newRid/></a:ext></a:extLst> into the last
         <a:blip> element (the picture we just inserted)
      Plus: write the SVG bytes to word/media/imageN.svg.
    """
    import re as _re
    import shutil
    import tempfile
    import zipfile
    from lxml import etree

    NS_REL = "{http://schemas.openxmlformats.org/package/2006/relationships}"
    NS_TYPES = "{http://schemas.openxmlformats.org/package/2006/content-types}"
    NS_A = "http://schemas.openxmlformats.org/drawingml/2006/main"
    NS_R = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
    NS_ASVG = "http://schemas.microsoft.com/office/drawing/2016/SVG/main"
    SVG_EXT_URI = "{96DAC541-7B7A-43D3-8B79-37D633B846F1}"

    svg_bytes = svg_path.read_bytes()

    # Read everything from the existing zip into memory; we'll rewrite it.
    with zipfile.ZipFile(docx_path, "r") as zin:
        names = zin.namelist()
        ct_data = zin.read("[Content_Types].xml")
        rels_data = zin.read("word/_rels/document.xml.rels")
        doc_data = zin.read("word/document.xml")
        kept: dict[str, bytes] = {
            n: zin.read(n) for n in names
            if n not in ("[Content_Types].xml",
                         "word/_rels/document.xml.rels",
                         "word/document.xml")
        }

    # 1. [Content_Types].xml — ensure <Default Extension="svg" ContentType="image/svg+xml"/>
    ct_root = etree.fromstring(ct_data)
    has_svg_default = any(
        el.tag.endswith("Default") and el.get("Extension", "").lower() == "svg"
        for el in ct_root
    )
    if not has_svg_default:
        default = etree.SubElement(ct_root, NS_TYPES + "Default")
        default.set("Extension", "svg")
        default.set("ContentType", "image/svg+xml")

    # 2. document.xml.rels — pick a fresh rId, fresh image filename
    rels_root = etree.fromstring(rels_data)
    existing_ids = [r.get("Id", "") for r in rels_root]
    max_n = 0
    for rid in existing_ids:
        m = _re.match(r"rId(\d+)$", rid)
        if m:
            max_n = max(max_n, int(m.group(1)))
    new_rid = f"rId{max_n + 1}"

    media_files = [n for n in kept if n.startswith("word/media/")]
    img_indices = []
    for n in media_files:
        m = _re.match(r"word/media/image(\d+)\.", n)
        if m:
            img_indices.append(int(m.group(1)))
    next_img = max(img_indices, default=0) + 1
    svg_filename = f"image{next_img}.svg"

    rel = etree.SubElement(rels_root, NS_REL + "Relationship")
    rel.set("Id", new_rid)
    rel.set(
        "Type",
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/image",
    )
    rel.set("Target", f"media/{svg_filename}")

    # 3. document.xml — patch the last <a:blip> with an <a:extLst><a:ext><asvg:svgBlip/></a:ext></a:extLst>
    doc_root = etree.fromstring(doc_data)
    blips = doc_root.findall(f".//{{{NS_A}}}blip")
    if not blips:
        raise RuntimeError(
            "_attach_svg_to_last_picture: no <a:blip> in document — "
            "was a picture actually inserted before this call?"
        )
    target_blip = blips[-1]
    extlst = etree.SubElement(target_blip, f"{{{NS_A}}}extLst")
    ext = etree.SubElement(extlst, f"{{{NS_A}}}ext")
    ext.set("uri", SVG_EXT_URI)
    svg_blip = etree.SubElement(
        ext, f"{{{NS_ASVG}}}svgBlip", nsmap={"asvg": NS_ASVG},
    )
    svg_blip.set(f"{{{NS_R}}}embed", new_rid)

    # Rewrite the zip atomically
    tmp_fd, tmp_name = tempfile.mkstemp(suffix=".docx", dir=str(docx_path.parent))
    import os
    os.close(tmp_fd)
    try:
        with zipfile.ZipFile(tmp_name, "w", zipfile.ZIP_DEFLATED) as zout:
            zout.writestr(
                "[Content_Types].xml",
                etree.tostring(ct_root, xml_declaration=True,
                               encoding="UTF-8", standalone=True),
            )
            zout.writestr(
                "word/_rels/document.xml.rels",
                etree.tostring(rels_root, xml_declaration=True,
                               encoding="UTF-8", standalone=True),
            )
            zout.writestr(
                "word/document.xml",
                etree.tostring(doc_root, xml_declaration=True,
                               encoding="UTF-8", standalone=True),
            )
            zout.writestr(f"word/media/{svg_filename}", svg_bytes)
            for name, data in kept.items():
                zout.writestr(name, data)
        shutil.move(tmp_name, str(docx_path))
    except Exception:
        try:
            os.unlink(tmp_name)
        except OSError:
            pass
        raise


# ============================================================
# PPTX edit primitives (python-pptx)
# ============================================================


def _pptx_iter_paragraphs(prs, include_notes: bool = True):
    """Yield every text paragraph in the deck (slide bodies + notes)."""
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    yield para
        if include_notes and slide.has_notes_slide:
            tf = slide.notes_slide.notes_text_frame
            for para in tf.paragraphs:
                yield para


def _pptx_set_paragraph_text(para, new_text: str) -> None:
    """Put new_text in the first run, blank the rest. Loses per-run formatting
    on edited text — same trade-off as DOCX find_replace."""
    if para.runs:
        para.runs[0].text = new_text
        for r in para.runs[1:]:
            r.text = ""
    else:
        # No runs (empty paragraph): add one
        para.add_run().text = new_text


def _pptx_find_replace(src: Path, dst: Path, pattern: str, replacement: str, regex: bool):
    from pptx import Presentation

    prs = Presentation(str(src))
    rx = re.compile(pattern) if regex else None
    for para in _pptx_iter_paragraphs(prs):
        if not para.runs:
            continue
        full = "".join(r.text for r in para.runs)
        new = rx.sub(replacement, full) if rx else full.replace(pattern, replacement)
        if new != full:
            _pptx_set_paragraph_text(para, new)
    prs.save(str(dst))


def _pptx_title(slide) -> Optional[str]:
    """Return the slide's title text, or None."""
    try:
        ts = slide.shapes.title
    except (AttributeError, KeyError):
        return None
    if ts is None or not ts.has_text_frame:
        return None
    return ts.text_frame.text


def _pptx_replace_heading(src: Path, dst: Path, old_title: str, new_title: str):
    from pptx import Presentation

    prs = Presentation(str(src))
    for slide in prs.slides:
        title = _pptx_title(slide)
        if title is None or title.strip() != old_title.strip():
            continue
        ts = slide.shapes.title
        paras = ts.text_frame.paragraphs
        if not paras:
            continue
        _pptx_set_paragraph_text(paras[0], new_title)
        for extra in paras[1:]:
            _pptx_set_paragraph_text(extra, "")
        prs.save(str(dst))
        return
    raise ValueError(f"slide title not found: {old_title!r}")


def _pptx_remove_slides(prs, indices_zero_based: list[int]):
    """Delete slides by zero-based index. Manipulates sldIdLst directly
    because python-pptx has no public delete_slide API."""
    sld_id_lst = prs.slides._sldIdLst  # noqa: SLF001
    items = list(sld_id_lst)
    for i in sorted(indices_zero_based, reverse=True):
        if 0 <= i < len(items):
            sld_id_lst.remove(items[i])


def _pptx_delete_slides_by_range(src: Path, dst: Path, page_range: str) -> int:
    from pptx import Presentation

    prs = Presentation(str(src))
    total = len(prs.slides)
    one_based = _parse_page_range(page_range, total)
    zero_based = [i - 1 for i in one_based]
    _pptx_remove_slides(prs, zero_based)
    prs.save(str(dst))
    return total - len(zero_based)


def _pptx_delete_section(src: Path, dst: Path, heading: str):
    """Delete every slide whose title equals `heading` (case-insensitive trim)."""
    from pptx import Presentation

    prs = Presentation(str(src))
    target = heading.strip().casefold()
    matches = [
        idx
        for idx, slide in enumerate(prs.slides)
        if (t := _pptx_title(slide)) and t.strip().casefold() == target
    ]
    if not matches:
        raise ValueError(f"no slide found with title: {heading!r}")
    _pptx_remove_slides(prs, matches)
    prs.save(str(dst))


def _pptx_insert_after(src: Path, dst: Path, anchor: str, content: str):
    """Insert `content` as a new bullet paragraph in the body of the slide
    whose title equals `anchor`. (For multi-line content, each line becomes
    its own paragraph.)"""
    from pptx import Presentation

    prs = Presentation(str(src))
    target = anchor.strip().casefold()
    for slide in prs.slides:
        title = _pptx_title(slide)
        if title is None or title.strip().casefold() != target:
            continue
        # Find the first non-title text frame to use as body
        body_tf = None
        title_shape = slide.shapes.title
        for shape in slide.shapes:
            if shape is title_shape:
                continue
            if shape.has_text_frame:
                body_tf = shape.text_frame
                break
        if body_tf is None:
            raise ValueError(
                f"slide '{anchor}' has no body placeholder; cannot insert"
            )
        for line in content.splitlines() or [content]:
            p = body_tf.add_paragraph()
            p.text = line
        prs.save(str(dst))
        return
    raise ValueError(f"slide title not found: {anchor!r}")


def _pptx_extract_text(src: Path) -> str:
    from pptx import Presentation

    prs = Presentation(str(src))
    chunks: list[str] = []
    for idx, slide in enumerate(prs.slides, start=1):
        title = _pptx_title(slide) or ""
        chunks.append(f"### Slide {idx}: {title.strip()}".rstrip(": "))
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = shape.text_frame.text.strip()
            if not txt or txt == title.strip():
                continue
            chunks.append(txt)
        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                chunks.append(f"[notes] {notes}")
    return "\n\n".join(chunks)


# ============================================================
# Public dispatcher
# ============================================================


async def find_replace(
    file_id: str,
    pattern: str,
    replacement: str,
    regex: bool = False,
) -> dict:
    src = _load_path(file_id)
    fmt = detect_from_path(src)
    out_id, dst = _new_output_for(src)
    t0 = time.time()

    matches = 0
    if fmt == "docx":
        matches = _docx_find_replace(src, dst, pattern, replacement, regex)
    elif fmt == "pptx":
        _pptx_find_replace(src, dst, pattern, replacement, regex)
        # python-pptx version doesn't track matches yet; approximate via diff
        matches = -1  # unknown
    elif fmt == "pdf":
        if regex:
            raise ValueError("PDF find_replace does not support regex (only literal text)")
        matches = _pdf_find_replace_inplace(src, dst, pattern, replacement)
    elif fmt in {"md", "html", "htm", "txt"}:
        text = src.read_text(encoding="utf-8")
        if regex:
            new_text, matches = re.subn(pattern, replacement, text)
        else:
            matches = text.count(pattern)
            new_text = text.replace(pattern, replacement) if matches else text
        dst.write_text(new_text, encoding="utf-8")
    else:
        raise ValueError(
            f"find_replace not supported for {fmt}. Convert to docx/md/html/pptx first."
        )

    return {
        "file_id": out_id,
        "url": public_url(out_id),
        "ms_elapsed": int((time.time() - t0) * 1000),
        "format": fmt,
        "matches": matches,
    }


async def delete_section(
    file_id: str, heading: str, level: Optional[int] = None
) -> dict:
    src = _load_path(file_id)
    fmt = detect_from_path(src)
    out_id, dst = _new_output_for(src)
    t0 = time.time()

    if fmt == "docx":
        _docx_delete_section(src, dst, heading, level)
    elif fmt == "pptx":
        _pptx_delete_section(src, dst, heading)
    elif fmt == "md":
        # delete from `# heading` (or matching level) until next heading of same/higher level
        lines = src.read_text(encoding="utf-8").splitlines()
        out_lines, deleting, found_lvl = [], False, None
        for ln in lines:
            m = re.match(r"^(#{1,6})\s+(.*)$", ln)
            if m:
                lvl = len(m.group(1))
                title = m.group(2).strip()
                if not deleting and title == heading.strip() and (level is None or lvl == level):
                    deleting = True
                    found_lvl = lvl
                    continue
                if deleting and lvl <= found_lvl:
                    deleting = False
            if not deleting:
                out_lines.append(ln)
        dst.write_text("\n".join(out_lines), encoding="utf-8")
    else:
        raise ValueError(f"delete_section not supported for {fmt}")

    return {
        "file_id": out_id,
        "url": public_url(out_id),
        "ms_elapsed": int((time.time() - t0) * 1000),
    }


async def insert_after(file_id: str, anchor: str, content: str) -> dict:
    src = _load_path(file_id)
    fmt = detect_from_path(src)
    out_id, dst = _new_output_for(src)
    t0 = time.time()

    if fmt == "docx":
        _docx_insert_after(src, dst, anchor, content)
    elif fmt == "pptx":
        _pptx_insert_after(src, dst, anchor, content)
    elif fmt == "md":
        lines = src.read_text(encoding="utf-8").splitlines()
        out_lines = []
        inserted = False
        for ln in lines:
            out_lines.append(ln)
            if not inserted and ln.strip() == anchor.strip():
                out_lines.append("")
                out_lines.extend(content.splitlines())
                inserted = True
            elif not inserted:
                m = re.match(r"^#{1,6}\s+(.*)$", ln)
                if m and m.group(1).strip() == anchor.strip():
                    out_lines.append("")
                    out_lines.extend(content.splitlines())
                    inserted = True
        if not inserted:
            raise ValueError(f"anchor not found: {anchor!r}")
        dst.write_text("\n".join(out_lines), encoding="utf-8")
    else:
        raise ValueError(f"insert_after not supported for {fmt}")

    return {
        "file_id": out_id,
        "url": public_url(out_id),
        "ms_elapsed": int((time.time() - t0) * 1000),
    }


async def replace_heading(file_id: str, old_title: str, new_title: str) -> dict:
    src = _load_path(file_id)
    fmt = detect_from_path(src)
    out_id, dst = _new_output_for(src)
    t0 = time.time()

    if fmt == "docx":
        _docx_replace_heading(src, dst, old_title, new_title)
    elif fmt == "pptx":
        _pptx_replace_heading(src, dst, old_title, new_title)
    elif fmt == "md":
        text = src.read_text(encoding="utf-8")
        pattern = re.compile(r"^(#{1,6}\s+)" + re.escape(old_title) + r"\s*$", re.M)
        new_text, n = pattern.subn(rf"\g<1>{new_title}", text)
        if n == 0:
            raise ValueError(f"heading not found: {old_title!r}")
        dst.write_text(new_text, encoding="utf-8")
    else:
        raise ValueError(f"replace_heading not supported for {fmt}")

    return {
        "file_id": out_id,
        "url": public_url(out_id),
        "ms_elapsed": int((time.time() - t0) * 1000),
    }


async def delete_pages(file_id: str, page_range: str) -> dict:
    """Delete pages (PDF) or slides (PPTX). page_range: '1-3', '5', '1,3,5-7'."""
    src = _load_path(file_id)
    fmt = detect_from_path(src)
    out_id, dst = _new_output_for(src)
    t0 = time.time()

    if fmt == "pdf":
        from pypdf import PdfReader, PdfWriter

        reader = PdfReader(str(src))
        total = len(reader.pages)
        to_remove = _parse_page_range(page_range, total)
        writer = PdfWriter()
        for i, page in enumerate(reader.pages):
            if (i + 1) not in to_remove:
                writer.add_page(page)
        with open(dst, "wb") as f:
            writer.write(f)
        remaining = total - len(to_remove)
    elif fmt == "pptx":
        remaining = _pptx_delete_slides_by_range(src, dst, page_range)
    else:
        raise ValueError(f"delete_pages only works on PDF/PPTX (got {fmt})")

    return {
        "file_id": out_id,
        "url": public_url(out_id),
        "pages_remaining": remaining,
        "ms_elapsed": int((time.time() - t0) * 1000),
        "format": fmt,
    }


def _parse_page_range(spec: str, total: int) -> set[int]:
    pages: set[int] = set()
    for part in spec.split(","):
        part = part.strip()
        if "-" in part:
            a, b = part.split("-", 1)
            pages.update(range(int(a), int(b) + 1))
        elif part:
            pages.add(int(part))
    return {p for p in pages if 1 <= p <= total}


async def merge(file_ids: list[str], target_format: str = "pdf") -> dict:
    """Merge multiple files into one (PDF only for now)."""
    if target_format != "pdf":
        raise ValueError("merge currently supports target_format=pdf only")

    from pypdf import PdfWriter

    t0 = time.time()
    out_id = new_file_id("pdf")
    dst = output_path(out_id)

    writer = PdfWriter()
    for fid in file_ids:
        p = _load_path(fid)
        if detect_from_path(p) != "pdf":
            raise ValueError(f"merge: all inputs must be PDF, got {p.suffix}")
        writer.append(str(p))
    with open(dst, "wb") as f:
        writer.write(f)

    return {
        "file_id": out_id,
        "url": public_url(out_id),
        "ms_elapsed": int((time.time() - t0) * 1000),
    }


async def docx_set_paragraph_format(
    file_id: str,
    anchor: str,
    *,
    page_break_before: Optional[bool] = None,
    keep_with_next: Optional[bool] = None,
    keep_lines_together: Optional[bool] = None,
    space_before_pt: Optional[float] = None,
    space_after_pt: Optional[float] = None,
) -> dict:
    """Apply DOCX paragraph-format flow controls to the paragraph whose text
    matches `anchor`. Use for layout micro-fixes:
       - page_break_before=True: force the heading to start on a new page
       - keep_with_next=True: never split this paragraph from the next one
         (e.g. heading should stay with the body that follows it)
       - keep_lines_together=True: don't break this paragraph mid-sentence"""
    from docx import Document
    from docx.shared import Pt

    src = _load_path(file_id)
    fmt = detect_from_path(src)
    if fmt != "docx":
        raise ValueError(f"docx_set_paragraph_format only works on DOCX (got {fmt})")
    out_id, dst = _new_output_for(src)
    t0 = time.time()

    doc = Document(str(src))
    target = anchor.strip()
    n = 0
    for p in doc.paragraphs:
        if p.text.strip() != target:
            continue
        pf = p.paragraph_format
        if page_break_before is not None:
            pf.page_break_before = page_break_before
        if keep_with_next is not None:
            pf.keep_with_next = keep_with_next
        if keep_lines_together is not None:
            pf.keep_together = keep_lines_together
        if space_before_pt is not None:
            pf.space_before = Pt(space_before_pt)
        if space_after_pt is not None:
            pf.space_after = Pt(space_after_pt)
        n += 1
    if n == 0:
        raise ValueError(f"no paragraph found with text: {anchor!r}")
    doc.save(str(dst))
    return {
        "file_id": out_id,
        "url": public_url(out_id),
        "paragraphs_updated": n,
        "ms_elapsed": int((time.time() - t0) * 1000),
    }


async def style_paragraph(file_id: str, anchor: str, **style) -> dict:
    """Apply text formatting (color, highlight, bold, italic, font_size, font_family)
    to every paragraph whose text equals `anchor`. DOCX only."""
    src = _load_path(file_id)
    fmt = detect_from_path(src)
    out_id, dst = _new_output_for(src)
    t0 = time.time()
    if fmt != "docx":
        raise ValueError(f"style_paragraph currently supports DOCX only (got {fmt})")
    n = _docx_style_paragraph(src, dst, anchor, **style)
    return {
        "file_id": out_id,
        "url": public_url(out_id),
        "ms_elapsed": int((time.time() - t0) * 1000),
        "paragraphs_styled": n,
    }


async def style_all_headings(file_id: str, levels: Optional[list[int]] = None,
                             **style) -> dict:
    """Apply text formatting to all heading paragraphs (optionally filtered to specific levels)."""
    src = _load_path(file_id)
    fmt = detect_from_path(src)
    out_id, dst = _new_output_for(src)
    t0 = time.time()
    if fmt != "docx":
        raise ValueError(f"style_all_headings currently supports DOCX only (got {fmt})")
    n = _docx_style_all_headings(src, dst, levels=levels, **style)
    return {
        "file_id": out_id,
        "url": public_url(out_id),
        "ms_elapsed": int((time.time() - t0) * 1000),
        "headings_styled": n,
    }


async def insert_image(file_id: str, anchor: str, image_file_id: str,
                       width_inches: Optional[float] = None,
                       svg_file_id: Optional[str] = None) -> dict:
    """Insert an image (already uploaded as a separate file_id) after the
    anchor paragraph in a DOCX.

    If `svg_file_id` is provided, the PNG is still embedded as a fallback
    but a vector SVG layer is attached on top (Microsoft's `<asvg:svgBlip>`
    extension). Modern Word (2016+) and LibreOffice 7+ render the SVG, so
    the image stays crisp at any zoom and in PDF export. Older readers
    silently fall back to the PNG. The caller typically passes both
    image_file_id and svg_file_id from a `chart_redraw(..., also_svg=true)`
    response.
    """
    src = _load_path(file_id)
    image_src = _load_path(image_file_id)
    fmt = detect_from_path(src)
    out_id, dst = _new_output_for(src)
    t0 = time.time()
    if fmt != "docx":
        raise ValueError(f"insert_image currently supports DOCX only (got {fmt})")
    _docx_insert_image_after(src, dst, anchor, image_src, width_inches=width_inches)
    vector_attached = False
    if svg_file_id:
        svg_src = _load_path(svg_file_id)
        _attach_svg_to_last_picture(dst, svg_src)
        vector_attached = True
    return {
        "file_id": out_id,
        "url": public_url(out_id),
        "ms_elapsed": int((time.time() - t0) * 1000),
        "vector_svg_attached": vector_attached,
    }


async def extract(file_id: str, mode: str = "text") -> dict:
    """Extract text/tables/metadata from a PDF/DOCX."""
    src = _load_path(file_id)
    fmt = detect_from_path(src)
    t0 = time.time()

    if mode == "text":
        if fmt == "pdf":
            import pdfplumber

            with pdfplumber.open(str(src)) as pdf:
                text = "\n\n".join((p.extract_text() or "") for p in pdf.pages)
            return {"text": text, "ms_elapsed": int((time.time() - t0) * 1000)}
        if fmt == "docx":
            from docx import Document

            doc = Document(str(src))
            text = "\n".join(p.text for p in doc.paragraphs)
            return {"text": text, "ms_elapsed": int((time.time() - t0) * 1000)}
        if fmt == "pptx":
            return {
                "text": _pptx_extract_text(src),
                "ms_elapsed": int((time.time() - t0) * 1000),
            }
        if fmt in {"md", "txt", "html", "htm"}:
            return {
                "text": src.read_text(encoding="utf-8"),
                "ms_elapsed": int((time.time() - t0) * 1000),
            }

    if mode == "tables" and fmt == "pdf":
        import pdfplumber

        tables = []
        with pdfplumber.open(str(src)) as pdf:
            for i, page in enumerate(pdf.pages, start=1):
                for t in page.extract_tables() or []:
                    tables.append({"page": i, "rows": t})
        return {"tables": tables, "ms_elapsed": int((time.time() - t0) * 1000)}

    raise ValueError(f"extract: unsupported mode={mode} for fmt={fmt}")
