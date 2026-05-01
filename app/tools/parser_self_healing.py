"""
Self-healing parser pipeline.

Tier-S #5. The user assessment: Kimi K2.6 is smarter than MinerU/Granite-Docling/
LlamaParse — so we don't actually integrate those. Instead we run a quality-aware
extraction:

  1. Primary pass: format-native parser (pdfplumber for PDF, python-docx for
     DOCX, etc) — fast, free, deterministic.
  2. Quality score: heuristics on the extracted text (char density per page,
     replacement-char ratio, printable-char ratio, whitespace ratio).
  3. If quality is below threshold (typical: scan-only PDF, broken CMap,
     embedded-font weirdness), fall back to Kimi K2.6 vision OCR per page.
  4. Return whichever method scored higher, with full provenance metadata.

Key insight: we don't try to detect "is this a scan?" upfront — we just
extract, score, and re-extract with vision if needed. This catches edge
cases like partial CMap corruption that other tools miss.

For documents that already parse cleanly, this adds only ~5ms of scoring
overhead. For broken PDFs, the vision fallback adds 1-3s per page but
recovers content that would otherwise be lost.
"""
from __future__ import annotations

import logging
import time
from pathlib import Path
from typing import Any, Optional

from app.storage import (
    input_path,
    new_file_id,
    output_path,
    public_url,
)
from app.tools.formats import detect_from_path

log = logging.getLogger("aice.parser_self_healing")


# ============================================================
# Quality scoring
# ============================================================


_REPLACEMENT_CHAR = "�"
_QUALITY_THRESHOLD = 0.55  # below this, try vision fallback


def score_text_quality(text: str, *, page_count: int = 1) -> dict[str, Any]:
    """Score extracted text on 0..1 scale based on heuristics.

    Penalties:
      - empty / very short → 0.0
      - high replacement-char (U+FFFD) ratio
      - low printable-char ratio
      - very low char density per page (likely scan with no text layer)
      - extreme whitespace ratio
    """
    if not text or len(text.strip()) < 10:
        return {
            "score": 0.0,
            "char_count": len(text),
            "page_count": page_count,
            "reason": "empty_or_too_short",
            "details": {},
        }

    n_chars = len(text)
    n_repl = text.count(_REPLACEMENT_CHAR)
    n_ws = sum(1 for c in text if c.isspace())
    n_printable = sum(1 for c in text if c.isprintable() or c in "\n\r\t")
    n_alpha = sum(1 for c in text if c.isalpha())

    repl_ratio = n_repl / n_chars
    ws_ratio = n_ws / n_chars
    printable_ratio = n_printable / n_chars
    alpha_ratio = n_alpha / n_chars
    chars_per_page = n_chars / max(page_count, 1)

    score = 1.0
    reasons = []

    # Density: scan-only PDFs typically yield <50 chars/page
    if chars_per_page < 50:
        score -= 0.5
        reasons.append("very_low_char_density")
    elif chars_per_page < 200:
        score -= 0.2
        reasons.append("low_char_density")

    # Garbled-encoding indicator
    if repl_ratio > 0.005:
        score -= min(0.5, repl_ratio * 50)
        reasons.append("replacement_chars")

    # Mostly-whitespace text → broken extraction
    if ws_ratio > 0.7:
        score -= 0.3
        reasons.append("excessive_whitespace")

    # Non-printable garbage
    if printable_ratio < 0.85:
        score -= 0.4
        reasons.append("non_printable_chars")

    # Almost no letters → likely just numbers/symbols, possibly bad
    if alpha_ratio < 0.3:
        score -= 0.2
        reasons.append("low_alpha_ratio")

    score = max(0.0, min(1.0, score))
    return {
        "score": round(score, 3),
        "char_count": n_chars,
        "page_count": page_count,
        "reason": ",".join(reasons) if reasons else "ok",
        "details": {
            "chars_per_page": round(chars_per_page, 1),
            "replacement_ratio": round(repl_ratio, 4),
            "whitespace_ratio": round(ws_ratio, 3),
            "printable_ratio": round(printable_ratio, 3),
            "alpha_ratio": round(alpha_ratio, 3),
        },
    }


# ============================================================
# Format-native primary extractors
# ============================================================


def _resolve(file_id: str) -> Path:
    p = output_path(file_id)
    if p.exists():
        return p
    p = input_path(file_id)
    if p.exists():
        return p
    raise FileNotFoundError(file_id)


def _extract_pdf_pdfplumber(src: Path) -> tuple[str, int]:
    """Returns (text, page_count)."""
    import pdfplumber
    pages = []
    with pdfplumber.open(str(src)) as pdf:
        for p in pdf.pages:
            pages.append(p.extract_text() or "")
        n = len(pdf.pages)
    return "\n\n".join(pages), n


def _extract_docx(src: Path) -> tuple[str, int]:
    from docx import Document
    doc = Document(str(src))
    text = "\n".join(p.text for p in doc.paragraphs)
    # Approximate page count: 35 paragraphs/page, minimum 1
    pages = max(1, len(doc.paragraphs) // 35)
    return text, pages


def _extract_pptx(src: Path) -> tuple[str, int]:
    from pptx import Presentation
    prs = Presentation(str(src))
    chunks = []
    for i, slide in enumerate(prs.slides, 1):
        chunks.append(f"### Slide {i}")
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    chunks.append(para.text)
    return "\n".join(chunks), len(prs.slides)


# ============================================================
# Kimi K2.6 vision fallback (per-page rasterize → OCR)
# ============================================================


async def _extract_pdf_via_vision(src: Path, *, dpi: int = 150,
                                   max_pages: Optional[int] = None) -> tuple[str, int]:
    """Rasterize each page → Kimi K2.6 OCR → join.

    Used as fallback when pdfplumber yields garbage. dpi=150 is the sweet spot:
    high enough for accurate OCR, low enough to keep token usage manageable.
    """
    import fitz
    from app.llm import vision as v

    doc = fitz.open(str(src))
    n = doc.page_count
    page_idxs = list(range(n if max_pages is None else min(n, max_pages)))
    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)

    page_texts: list[str] = []
    for i in page_idxs:
        pix = doc[i].get_pixmap(matrix=matrix)
        png_fid = new_file_id("png")
        png_path = output_path(png_fid)
        with open(png_path, "wb") as f:
            f.write(pix.tobytes("png"))
        try:
            text = await v.ocr_image(png_path)
        except Exception as e:
            log.warning("vision OCR failed on page %d: %s", i + 1, e)
            text = ""
        page_texts.append(text)
    doc.close()

    return "\n\n".join(page_texts), n


async def _extract_image_via_vision(src: Path) -> tuple[str, int]:
    from app.llm import vision as v
    text = await v.ocr_image(src)
    return text, 1


# ============================================================
# Public API
# ============================================================


async def parse_with_quality(
    file_id: str,
    *,
    force_vision: bool = False,
    quality_threshold: float = _QUALITY_THRESHOLD,
    max_vision_pages: Optional[int] = None,
) -> dict[str, Any]:
    """Self-healing extract: native parser → quality score → vision fallback
    if score < threshold. Returns whichever method scored higher.

    Args:
      file_id: source document
      force_vision: skip native parser, go straight to Kimi vision
      quality_threshold: minimum acceptable native-parser score (default 0.55)
      max_vision_pages: cap vision OCR to N pages (None = all)

    Returns:
      {
        "text": str,
        "method": "pdfplumber" | "kimi_vision" | "docx_native" | "pptx_native" | ...,
        "quality_score": float,
        "page_count": int,
        "fallbacks_tried": list[dict],   # method + score for each attempt
        "ms_elapsed": int,
      }
    """
    src = _resolve(file_id)
    fmt = detect_from_path(src)
    t0 = time.time()
    fallbacks: list[dict[str, Any]] = []

    # Native primary pass (skipped if force_vision=True)
    primary_text, primary_pages, primary_method = "", 0, ""
    primary_score = None

    if not force_vision:
        if fmt == "pdf":
            primary_method = "pdfplumber"
            primary_text, primary_pages = _extract_pdf_pdfplumber(src)
        elif fmt == "docx":
            primary_method = "docx_native"
            primary_text, primary_pages = _extract_docx(src)
        elif fmt == "pptx":
            primary_method = "pptx_native"
            primary_text, primary_pages = _extract_pptx(src)
        elif fmt in {"png", "jpg", "jpeg", "webp", "gif", "bmp"}:
            # Images have no native text — go straight to vision
            primary_method = "kimi_vision"
            primary_text, primary_pages = await _extract_image_via_vision(src)
        elif fmt in {"md", "txt", "html", "htm"}:
            primary_method = "text_native"
            primary_text = src.read_text(encoding="utf-8")
            primary_pages = 1
        else:
            raise ValueError(f"unsupported format for self-healing parser: {fmt}")

        primary_q = score_text_quality(primary_text, page_count=primary_pages)
        primary_score = primary_q["score"]
        fallbacks.append({
            "method": primary_method, "score": primary_score,
            "char_count": primary_q["char_count"],
            "reason": primary_q["reason"],
        })

        # Good enough? Return early.
        if primary_score >= quality_threshold:
            return {
                "text": primary_text,
                "method": primary_method,
                "quality_score": primary_score,
                "quality_details": primary_q["details"],
                "page_count": primary_pages,
                "fallbacks_tried": fallbacks,
                "ms_elapsed": int((time.time() - t0) * 1000),
            }

    # Vision fallback (only meaningful for PDF or images)
    if fmt == "pdf" or force_vision:
        log.info(
            "parse_with_quality: primary_score=%s < %s — trying vision fallback",
            primary_score, quality_threshold,
        )
        if fmt == "pdf":
            vision_text, vision_pages = await _extract_pdf_via_vision(
                src, max_pages=max_vision_pages,
            )
        elif fmt in {"png", "jpg", "jpeg", "webp", "gif", "bmp"}:
            vision_text, vision_pages = await _extract_image_via_vision(src)
        else:
            # No vision option for this format — return native pass
            return {
                "text": primary_text, "method": primary_method,
                "quality_score": primary_score,
                "page_count": primary_pages,
                "fallbacks_tried": fallbacks,
                "ms_elapsed": int((time.time() - t0) * 1000),
                "warning": "primary_score below threshold but no vision fallback for format",
            }

        vision_q = score_text_quality(vision_text, page_count=vision_pages)
        fallbacks.append({
            "method": "kimi_vision", "score": vision_q["score"],
            "char_count": vision_q["char_count"], "reason": vision_q["reason"],
        })

        # Pick the better of the two
        if force_vision or vision_q["score"] > (primary_score or 0):
            return {
                "text": vision_text,
                "method": "kimi_vision",
                "quality_score": vision_q["score"],
                "quality_details": vision_q["details"],
                "page_count": vision_pages,
                "fallbacks_tried": fallbacks,
                "ms_elapsed": int((time.time() - t0) * 1000),
            }

    # Vision didn't beat primary (or wasn't applicable) — return primary
    return {
        "text": primary_text,
        "method": primary_method,
        "quality_score": primary_score,
        "page_count": primary_pages,
        "fallbacks_tried": fallbacks,
        "ms_elapsed": int((time.time() - t0) * 1000),
    }
