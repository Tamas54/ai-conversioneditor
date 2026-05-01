"""
PPTX builder — agentic slide construction tools.

Unlike `convert_format(target=pptx)` which mechanically dumps content into a
single fake slide, these tools let the agent BUILD a presentation slide by slide:
  1. pptx_create — start an empty deck with a title slide
  2. pptx_add_slide — add a slide with title + body (free text or multi-line)
  3. pptx_add_bullets_slide — add a slide with title + bullet list
  4. pptx_set_slide_styles — apply font/color/size to one or all slides
  5. pptx_slide_count — read how many slides the deck has

The agent reads the source document, plans the slide structure itself
(creative, content-aware), and calls these tools to build it.
"""
from __future__ import annotations

import logging
import time
from copy import deepcopy
from pathlib import Path
from typing import Any, Optional

from app.storage import (
    input_path,
    new_file_id,
    output_path,
    public_url,
)

log = logging.getLogger("aice.pptx_builder")


# Layout indices (default python-pptx layouts):
#   0 = Title Slide (title + subtitle)
#   1 = Title and Content (title + body placeholder)
#   2 = Section Header (large title)
#   3 = Two Content (title + two body columns)
#   5 = Title Only
LAYOUT_NAMES = {
    "title": 0,
    "title_content": 1,
    "section_header": 2,
    "two_column": 3,
    "title_only": 5,
}


def _resolve(file_id: str) -> Path:
    p = output_path(file_id)
    if p.exists():
        return p
    p = input_path(file_id)
    if p.exists():
        return p
    raise FileNotFoundError(f"file_id not found: {file_id}")


def _new_pptx_output() -> tuple[str, Path]:
    fid = new_file_id("pptx")
    return fid, output_path(fid)


def _parse_rgb(color: str) -> tuple[int, int, int]:
    """Reuse the same color parsing as docx styling (hex or named)."""
    named = {
        "black": "000000", "white": "FFFFFF", "red": "C0392B", "green": "27AE60",
        "blue": "2C3E50", "navy": "1F3A5F", "gold": "B8860B", "orange": "E67E22",
        "purple": "8E44AD", "teal": "16A085", "gray": "7F8C8D", "darkgray": "555555",
        "lightgray": "BDBDBD",
        "darkblue": "1F3A5F", "darkgreen": "1B5E20", "darkred": "8B1F1F",
    }
    s = color.strip().lstrip("#")
    if s.lower() in named:
        s = named[s.lower()]
    if len(s) != 6:
        raise ValueError(f"invalid color: {color!r}")
    return (int(s[0:2], 16), int(s[2:4], 16), int(s[4:6], 16))


# ============================================================
# Slide-building primitives
# ============================================================


async def pptx_plan(
    deck_title: str,
    slides: list[dict[str, Any]],
    deck_subtitle: Optional[str] = None,
) -> dict[str, Any]:
    """Record the planned slide structure BEFORE building anything. This is
    a no-op step that just echoes the plan back — its real purpose is to
    force the agent to commit to a structure before adding slides.

    `slides`: list of {title, content_type, key_points/headers/rows} dicts.
      content_type ∈ {'bullets', 'content', 'table'}
        - 'bullets' → will be built with pptx_add_bullets_slide; key_points = list[str]
        - 'content' → will be built with pptx_add_slide (free-form body); key_points = str
        - 'table'   → will be built with pptx_add_table_slide; provide headers + rows
    """
    summary = []
    for i, s in enumerate(slides, start=1):
        ct = s.get("content_type", "bullets")
        title = s.get("title", f"Slide {i}")
        if ct == "table":
            shape = f"table ({len(s.get('headers', []))}c × {len(s.get('rows', []))}r)"
        elif ct == "bullets":
            shape = f"{len(s.get('key_points', []))} bullets"
        else:
            kp = s.get("key_points", "")
            shape = f"content ({len(kp.splitlines() if isinstance(kp,str) else kp)} lines)"
        summary.append(f"#{i} [{ct}] {title} — {shape}")
    return {
        "deck_title": deck_title,
        "deck_subtitle": deck_subtitle,
        "planned_slide_count": len(slides),
        "summary": summary,
        "note": (
            "Plan recorded. Next: call pptx_create(title=deck_title, subtitle=deck_subtitle), "
            "then build each slide one-by-one according to the plan, choosing "
            "pptx_add_bullets_slide / pptx_add_slide / pptx_add_table_slide based on content_type. "
            "After each add, optionally call pptx_set_slide_styles(slide_index=last_index, ...) "
            "to fine-tune that single slide. Final deck-wide styling at the end."
        ),
    }


async def pptx_create(title: str, subtitle: Optional[str] = None) -> dict[str, Any]:
    """Create a brand-new presentation. The first slide is a Title slide
    with the given title and optional subtitle. Returns the new file_id."""
    from pptx import Presentation

    t0 = time.time()
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_NAMES["title"]])
    _set_title_text(slide, title)
    if subtitle and len(slide.placeholders) > 1:
        try:
            sub_shape = slide.placeholders[1]
            if sub_shape.has_text_frame:
                _set_body_text(sub_shape, subtitle, default_size_pt=18)
        except KeyError:
            pass

    fid, dst = _new_pptx_output()
    prs.save(str(dst))
    return {
        "file_id": fid,
        "url": public_url(fid),
        "slides": len(prs.slides),
        "ms_elapsed": int((time.time() - t0) * 1000),
    }


def _set_body_text(
    body_shape, body: str, *,
    default_size_pt: float = 16.0,
    no_autofit: bool = True,
    paragraph_space_after_pt: float = 8.0,
) -> None:
    """Put `body` into a body shape with predictable typography:
    - vertical_anchor = TOP (text starts from top, never bottom-crushed)
    - word_wrap = True (lines wrap inside the shape)
    - auto_size = NONE by default (avoids the python-pptx aggressive shrink
      that makes text 'ridiculously small' on dense slides). Pair this with
      explicit textbox geometry and content-density limits in the prompt.
    - Each line in `body` becomes its own paragraph with `space_after` to
      visually separate ideas (works as soft bullets without unicode chars)."""
    from pptx.util import Pt
    from pptx.enum.text import MSO_ANCHOR, MSO_AUTO_SIZE

    tf = body_shape.text_frame
    tf.word_wrap = True
    try:
        tf.auto_size = MSO_AUTO_SIZE.NONE if no_autofit else MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    except (ValueError, NotImplementedError):
        pass
    try:
        tf.vertical_anchor = MSO_ANCHOR.TOP
    except (ValueError, NotImplementedError):
        pass
    # Tighter internal margins so we use the full text box area
    tf.margin_left = tf.margin_right = Pt(6)
    tf.margin_top = tf.margin_bottom = Pt(4)

    lines = body.splitlines() if body else [""]
    if not lines:
        lines = [""]

    # First paragraph
    tf.text = lines[0]
    p0 = tf.paragraphs[0]
    p0.space_after = Pt(paragraph_space_after_pt)
    for run in p0.runs or [p0.add_run()]:
        run.font.size = Pt(default_size_pt)

    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = line
        p.space_after = Pt(paragraph_space_after_pt)
        for run in p.runs or [p.add_run()]:
            run.font.size = Pt(default_size_pt)


def _set_title_text(slide, title: str) -> None:
    """Set the slide title with auto-shrink for long titles to prevent
    line-wrap overflow against the title placeholder."""
    from pptx.util import Pt

    if not slide.shapes.title:
        return
    slide.shapes.title.text_frame.word_wrap = True
    slide.shapes.title.text = title
    # Auto-shrink: long titles get smaller font so they fit on 1-2 lines
    if len(title) > 60:
        sz = 24
    elif len(title) > 40:
        sz = 28
    else:
        sz = 32
    for para in slide.shapes.title.text_frame.paragraphs:
        for run in para.runs or []:
            run.font.size = Pt(sz)


def _find_body_placeholder(slide):
    """Return the first non-title placeholder that can hold body text, or None."""
    title_shape = slide.shapes.title
    for shape in slide.placeholders:
        if shape is title_shape:
            continue
        if shape.has_text_frame:
            return shape
    return None


async def pptx_add_slide(
    file_id: str,
    title: str,
    body: str = "",
    layout: str = "title_content",
) -> dict[str, Any]:
    """Add a slide to an existing presentation.
    - layout: 'title' | 'title_content' | 'section_header' | 'two_column' | 'title_only'
    - body: free text; '\\n' in the body becomes a new paragraph (bullet) on the slide.
    Returns the new file_id and the new total slide count."""
    from pptx import Presentation

    src = _resolve(file_id)
    t0 = time.time()
    prs = Presentation(str(src))
    layout_idx = LAYOUT_NAMES.get(layout, LAYOUT_NAMES["title_content"])
    if layout_idx >= len(prs.slide_layouts):
        layout_idx = LAYOUT_NAMES["title_content"]
    # 'section_header' uses the layout's centered title (no body); other
    # content slides use title_only + a manually-placed textbox so we control
    # geometry precisely and avoid the python-pptx blank-theme placeholder
    # quirks (random vertical anchors, oversized text, top/bottom crushing).
    from pptx.util import Inches

    if layout == "section_header":
        slide = prs.slides.add_slide(prs.slide_layouts[layout_idx])
        _set_title_text(slide, title)
        body_shape = _find_body_placeholder(slide)
        if body and body_shape is not None:
            _set_body_text(body_shape, body, default_size_pt=18, no_autofit=True)
    else:
        # Always use title_only (5) + explicit textbox for predictable layout
        slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_NAMES["title_only"]])
        _set_title_text(slide, title)
        if body:
            tb = slide.shapes.add_textbox(
                Inches(0.6),     # left margin
                Inches(2.0),     # top: below brand chrome (eyebrow + section label + title + rule)
                Inches(8.8),
                Inches(5.0),     # leaves ~0.4" for footer / page numbers
            )
            _set_body_text(tb, body, default_size_pt=16, no_autofit=True)

    fid, dst = _new_pptx_output()
    prs.save(str(dst))
    return {
        "file_id": fid,
        "url": public_url(fid),
        "slides": len(prs.slides),
        "added_slide_index": len(prs.slides) - 1,
        "ms_elapsed": int((time.time() - t0) * 1000),
    }


async def pptx_add_bullets_slide(
    file_id: str,
    title: str,
    bullets: list[str],
    layout: str = "title_content",
) -> dict[str, Any]:
    """Convenience: add a slide with title + bullet list. Each bullet is its own paragraph."""
    body = "\n".join(b.lstrip("•- *").strip() for b in (bullets or []))
    return await pptx_add_slide(file_id, title, body=body, layout=layout)


async def pptx_add_table_slide(
    file_id: str,
    title: str,
    headers: list[str],
    rows: list[list[str]],
    *,
    header_color: str = "#1F3A5F",
    header_text_color: str = "#FFFFFF",
    row_band_color: str = "#F4E7C5",
) -> dict[str, Any]:
    """Add a slide with a real PPTX TABLE shape (not a body-text dump).
    Use this when the content is genuinely tabular (compare values across
    categories). Auto-styled header row + alternating row bands."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    src = _resolve(file_id)
    t0 = time.time()
    prs = Presentation(str(src))
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_NAMES["title_only"]])
    _set_title_text(slide, title)

    n_cols = max(1, len(headers))
    n_rows = max(2, len(rows) + 1)  # header + at least 1 data row

    # Slide is 10" × 7.5" by default; place table below brand chrome
    left = Inches(0.5)
    top = Inches(2.0)
    width = Inches(9)
    # Height scales with rows: ~0.4" per row, capped to 5.0" (leaves footer room)
    height = Inches(min(5.0, 0.4 * n_rows + 0.3))

    table_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
    table = table_shape.table

    # Proportional column widths: longer header label or longer data → wider column.
    # Compute the max content length per column from header + first 4 data rows.
    col_lens: list[int] = []
    for ci in range(n_cols):
        header_len = len(str(headers[ci]) if ci < len(headers) else "")
        sample_max = 0
        for r in rows[:4]:
            if isinstance(r, (list, tuple)) and ci < len(r):
                sample_max = max(sample_max, len(str(r[ci])))
        col_lens.append(max(6, max(header_len, sample_max)))
    total = sum(col_lens) or 1
    width_emu = width
    if hasattr(width_emu, "emu"):
        width_emu = width_emu.emu
    for ci, col in enumerate(table.columns):
        col.width = int(width_emu * col_lens[ci] / total)

    header_rgb = RGBColor(*_parse_rgb(header_color))
    header_text_rgb = RGBColor(*_parse_rgb(header_text_color))
    band_rgb = RGBColor(*_parse_rgb(row_band_color))

    # Header row — readable size, tight margins so the text uses the cell width
    for col_idx, h in enumerate(headers[:n_cols]):
        cell = table.cell(0, col_idx)
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_rgb
        cell.text = str(h)
        cell.text_frame.word_wrap = True
        cell.margin_left = cell.margin_right = Pt(4)
        cell.margin_top = cell.margin_bottom = Pt(3)
        for para in cell.text_frame.paragraphs:
            for run in (para.runs or [para.add_run()]):
                run.font.bold = True
                run.font.size = Pt(12)
                run.font.color.rgb = header_text_rgb

    # Data rows
    for row_idx, row in enumerate(rows, start=1):
        if row_idx >= n_rows:
            break
        if not isinstance(row, (list, tuple)):
            row = [row]
        # Alternating band on odd-indexed data rows
        if row_idx % 2 == 0:
            for col_idx in range(n_cols):
                c = table.cell(row_idx, col_idx)
                c.fill.solid()
                c.fill.fore_color.rgb = band_rgb
        for col_idx in range(n_cols):
            cell = table.cell(row_idx, col_idx)
            val = row[col_idx] if col_idx < len(row) else ""
            cell.text = str(val) if val is not None else ""
            cell.text_frame.word_wrap = True
            cell.margin_left = cell.margin_right = Pt(4)
            cell.margin_top = cell.margin_bottom = Pt(3)
            for para in cell.text_frame.paragraphs:
                for run in (para.runs or [para.add_run()]):
                    run.font.size = Pt(11)

    fid, dst = _new_pptx_output()
    prs.save(str(dst))
    return {
        "file_id": fid,
        "url": public_url(fid),
        "slides": len(prs.slides),
        "added_slide_index": len(prs.slides) - 1,
        "table_rows": n_rows,
        "table_cols": n_cols,
        "ms_elapsed": int((time.time() - t0) * 1000),
    }


async def pptx_set_slide_styles(
    file_id: str,
    slide_index: Optional[int] = None,
    *,
    title_color: Optional[str] = None,
    body_color: Optional[str] = None,
    title_font_size: Optional[float] = None,
    body_font_size: Optional[float] = None,
    font_family: Optional[str] = None,
    title_bold: Optional[bool] = None,
    title_italic: Optional[bool] = None,
) -> dict[str, Any]:
    """Apply visual styling to one slide (slide_index, 0-based) or to all slides
    (slide_index=None). For text already in the slide; doesn't change layouts.

    All style parameters are optional — only the ones you pass are applied."""
    from pptx import Presentation
    from pptx.util import Pt
    from pptx.dml.color import RGBColor

    src = _resolve(file_id)
    t0 = time.time()
    prs = Presentation(str(src))

    title_rgb = RGBColor(*_parse_rgb(title_color)) if title_color else None
    body_rgb = RGBColor(*_parse_rgb(body_color)) if body_color else None

    if slide_index is None:
        slides = list(prs.slides)
    else:
        if slide_index < 0 or slide_index >= len(prs.slides):
            raise ValueError(f"slide_index {slide_index} out of range (0..{len(prs.slides) - 1})")
        slides = [prs.slides[slide_index]]

    def _apply_run(run, *, color=None, size=None, family=None,
                   bold=None, italic=None):
        if color is not None:
            run.font.color.rgb = color
        if size is not None:
            run.font.size = Pt(size)
        if family:
            run.font.name = family
        if bold is not None:
            run.font.bold = bold
        if italic is not None:
            run.font.italic = italic

    styled = 0
    for slide in slides:
        # Title shape
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in para.runs or [para.add_run()]:
                    _apply_run(
                        run, color=title_rgb, size=title_font_size,
                        family=font_family, bold=title_bold, italic=title_italic,
                    )
        # Body shapes
        for shape in slide.placeholders:
            if shape is slide.shapes.title:
                continue
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in para.runs or [para.add_run()]:
                    _apply_run(
                        run, color=body_rgb, size=body_font_size, family=font_family,
                    )
        # Free-floating textboxes added by us
        for shape in slide.shapes:
            if shape.has_text_frame and shape not in list(slide.placeholders):
                for para in shape.text_frame.paragraphs:
                    for run in para.runs or [para.add_run()]:
                        _apply_run(
                            run, color=body_rgb, size=body_font_size, family=font_family,
                        )
        styled += 1

    fid, dst = _new_pptx_output()
    prs.save(str(dst))
    return {
        "file_id": fid,
        "url": public_url(fid),
        "slides_styled": styled,
        "ms_elapsed": int((time.time() - t0) * 1000),
    }


async def pptx_add_chart_slide(
    file_id: str,
    title: str,
    chart_type: str = "column",
    categories: Optional[list[str]] = None,
    series: Optional[list[dict[str, Any]]] = None,
    *,
    legend: bool = True,
    legend_position: str = "bottom",      # 'top' | 'bottom' | 'left' | 'right'
    palette: Optional[list[str]] = None,  # list of hex colors for series
    title_axis_x: Optional[str] = None,
    title_axis_y: Optional[str] = None,
    show_data_labels: bool = False,
) -> dict[str, Any]:
    """Add a slide with a NATIVE PPTX chart (editable in PowerPoint/LibreOffice).

    chart_type: 'column' (vertical bars, default), 'bar' (horizontal),
                'line', 'pie', 'doughnut', 'area', 'scatter'
    categories: x-axis category labels, e.g. ['S1', 'S2', 'S3', 'S4']
    series: [{'name': 'Brent (USD)', 'values': [75, 93, 120, 155]}, ...]
    palette: optional list of hex colors to use for the series
             (e.g. ['#1F3A5F', '#C66A00'] for HNI brand)
    """
    from pptx import Presentation
    from pptx.chart.data import CategoryChartData, XyChartData
    from pptx.dml.color import RGBColor
    from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
    from pptx.util import Inches, Pt

    src = _resolve(file_id)
    t0 = time.time()
    prs = Presentation(str(src))
    slide = prs.slides.add_slide(prs.slide_layouts[LAYOUT_NAMES["title_only"]])
    _set_title_text(slide, title)

    # Prepare chart data
    chart_type_map = {
        "column": XL_CHART_TYPE.COLUMN_CLUSTERED,
        "column_stacked": XL_CHART_TYPE.COLUMN_STACKED,
        "bar": XL_CHART_TYPE.BAR_CLUSTERED,
        "bar_stacked": XL_CHART_TYPE.BAR_STACKED,
        "line": XL_CHART_TYPE.LINE,
        "line_markers": XL_CHART_TYPE.LINE_MARKERS,
        "pie": XL_CHART_TYPE.PIE,
        "doughnut": XL_CHART_TYPE.DOUGHNUT,
        "area": XL_CHART_TYPE.AREA,
        "area_stacked": XL_CHART_TYPE.AREA_STACKED,
        "scatter": XL_CHART_TYPE.XY_SCATTER,
        "radar": XL_CHART_TYPE.RADAR,
    }
    xl_type = chart_type_map.get(chart_type, XL_CHART_TYPE.COLUMN_CLUSTERED)

    chart_data = CategoryChartData()
    chart_data.categories = categories or []
    for s in series or []:
        chart_data.add_series(s.get("name", "Series"), s.get("values", []))

    # Position chart in the content area (below brand chrome)
    left = Inches(0.6)
    top = Inches(2.0)
    width = Inches(8.8)
    height = Inches(4.8)

    chart_shape = slide.shapes.add_chart(xl_type, left, top, width, height, chart_data)
    chart = chart_shape.chart

    # Legend
    if legend and len(series or []) > 0:
        chart.has_legend = True
        legend_pos_map = {
            "top": XL_LEGEND_POSITION.TOP,
            "bottom": XL_LEGEND_POSITION.BOTTOM,
            "left": XL_LEGEND_POSITION.LEFT,
            "right": XL_LEGEND_POSITION.RIGHT,
        }
        try:
            chart.legend.position = legend_pos_map.get(legend_position, XL_LEGEND_POSITION.BOTTOM)
            chart.legend.include_in_layout = False
            chart.legend.font.size = Pt(10)
        except Exception:
            pass
    else:
        chart.has_legend = False

    # Apply palette to series fills
    if palette and chart_type not in {"scatter"}:
        try:
            for plot in chart.plots:
                for i, ser in enumerate(plot.series):
                    color = palette[i % len(palette)]
                    rgb = RGBColor(*_parse_rgb(color))
                    if chart_type in {"line", "line_markers"}:
                        ser.format.line.color.rgb = rgb
                        ser.format.line.width = Pt(2.5)
                    else:
                        ser.format.fill.solid()
                        ser.format.fill.fore_color.rgb = rgb
        except Exception as e:
            log.debug("palette apply failed: %s", e)

    # Axis titles
    if title_axis_x or title_axis_y:
        try:
            if title_axis_x:
                chart.category_axis.has_title = True
                chart.category_axis.axis_title.text_frame.text = title_axis_x
            if title_axis_y:
                chart.value_axis.has_title = True
                chart.value_axis.axis_title.text_frame.text = title_axis_y
        except Exception:
            pass

    # Data labels
    if show_data_labels:
        try:
            for plot in chart.plots:
                plot.has_data_labels = True
                plot.data_labels.font.size = Pt(9)
        except Exception:
            pass

    # Disable chart's own title (use slide title instead)
    chart.has_title = False

    fid, dst = _new_pptx_output()
    prs.save(str(dst))
    return {
        "file_id": fid,
        "url": public_url(fid),
        "slides": len(prs.slides),
        "added_slide_index": len(prs.slides) - 1,
        "chart_type": chart_type,
        "n_series": len(series or []),
        "n_categories": len(categories or []),
        "ms_elapsed": int((time.time() - t0) * 1000),
    }


async def pptx_apply_theme(
    file_id: str,
    *,
    accent_color: str = "#1F3A5F",
    accent_position: str = "left",        # 'left' | 'top' | 'none'
    accent_thickness_in: float = 0.18,
    footer_text: Optional[str] = None,
    footer_color: str = "#7F6A2C",
    page_numbers: bool = True,
    skip_first_slide: bool = True,        # don't decorate the cover
    title_color: Optional[str] = "#1F3A5F",
    title_bold: bool = True,
    body_color: Optional[str] = "#2E1F06",
    font_family: Optional[str] = "Calibri",
) -> dict[str, Any]:
    """Apply a unified visual theme to ALL slides: accent bar + footer +
    page numbers + cohesive title/body styling. Cover slide (slide 0) gets
    no chrome by default — keep it clean as the hero.

    Position 'left' = thin vertical bar on the left edge (good for content
    decks). 'top' = thin horizontal bar at top (good for short reports).
    'none' = skip the bar but keep footer/page numbers."""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.dml.color import RGBColor

    src = _resolve(file_id)
    t0 = time.time()
    prs = Presentation(str(src))
    sw = prs.slide_width
    sh = prs.slide_height

    accent_rgb = RGBColor(*_parse_rgb(accent_color)) if accent_color else None
    footer_rgb = RGBColor(*_parse_rgb(footer_color))
    title_rgb = RGBColor(*_parse_rgb(title_color)) if title_color else None
    body_rgb = RGBColor(*_parse_rgb(body_color)) if body_color else None

    decorated = 0
    for idx, slide in enumerate(prs.slides):
        skip_chrome = skip_first_slide and idx == 0

        # Apply title/body styling on every slide (including cover) for cohesion
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            for para in slide.shapes.title.text_frame.paragraphs:
                for run in (para.runs or [para.add_run()]):
                    if title_rgb is not None:
                        run.font.color.rgb = title_rgb
                    if title_bold is not None:
                        run.font.bold = title_bold
                    if font_family:
                        run.font.name = font_family
        for shape in slide.placeholders:
            if shape is slide.shapes.title or not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                for run in (para.runs or [para.add_run()]):
                    if body_rgb is not None:
                        run.font.color.rgb = body_rgb
                    if font_family:
                        run.font.name = font_family

        if skip_chrome:
            continue

        # Accent bar
        if accent_rgb is not None and accent_position != "none":
            thickness_emu = int(Inches(accent_thickness_in))
            if accent_position == "left":
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, 0, thickness_emu, sh,
                )
            else:  # top
                bar = slide.shapes.add_shape(
                    MSO_SHAPE.RECTANGLE, 0, 0, sw, thickness_emu,
                )
            bar.fill.solid()
            bar.fill.fore_color.rgb = accent_rgb
            bar.line.fill.background()  # remove border outline
            # Move to back so it doesn't cover content
            spTree = bar._element.getparent()
            spTree.remove(bar._element)
            spTree.insert(2, bar._element)  # behind shapes (after nvGrpSpPr+grpSpPr)

        # Footer & page number row
        footer_top = sh - Inches(0.35)
        if footer_text:
            tb = slide.shapes.add_textbox(Inches(0.4), footer_top, sw - Inches(1.0), Inches(0.3))
            tb.text_frame.text = footer_text
            for run in tb.text_frame.paragraphs[0].runs or [tb.text_frame.paragraphs[0].add_run()]:
                run.font.size = Pt(9)
                run.font.color.rgb = footer_rgb
                if font_family:
                    run.font.name = font_family
        if page_numbers:
            pn = slide.shapes.add_textbox(sw - Inches(0.7), footer_top, Inches(0.5), Inches(0.3))
            pn.text_frame.text = f"{idx + 1} / {len(prs.slides)}"
            for run in pn.text_frame.paragraphs[0].runs or [pn.text_frame.paragraphs[0].add_run()]:
                run.font.size = Pt(9)
                run.font.color.rgb = footer_rgb
                if font_family:
                    run.font.name = font_family
            pn.text_frame.paragraphs[0].alignment = 2  # right-align (PP_ALIGN.RIGHT)

        decorated += 1

    fid, dst = _new_pptx_output()
    prs.save(str(dst))
    return {
        "file_id": fid,
        "url": public_url(fid),
        "slides_decorated": decorated,
        "ms_elapsed": int((time.time() - t0) * 1000),
    }


async def pptx_modify_textbox(
    file_id: str,
    find_text: str,
    *,
    slide_index: Optional[int] = None,
    new_text: Optional[str] = None,
    font_size_pt: Optional[float] = None,
    bold: Optional[bool] = None,
    italic: Optional[bool] = None,
    font_color: Optional[str] = None,
    left_inches: Optional[float] = None,
    top_inches: Optional[float] = None,
    width_inches: Optional[float] = None,
    height_inches: Optional[float] = None,
) -> dict[str, Any]:
    """Find a textbox/shape by partial text match and modify its style and/or
    geometry. Critical for self-correcting visual builds: when a render shows
    text overlap or sizing problems, the agent uses this to fix specific
    textboxes WITHOUT rebuilding the whole slide.

    Args:
      find_text: substring to match against shape text (case-sensitive).
        First matching shape across the slide(s) is modified.
      slide_index: limit search to one slide; None = search all slides
      new_text: replace the entire shape text content
      font_size_pt: change all runs to this size
      bold/italic: change emphasis on all runs
      font_color: hex like '#FFFFFF'
      left/top/width/height_inches: reposition + resize the shape

    Returns: file_id, modified_slide_index, shape_index, what changed.
    """
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor

    src = _resolve(file_id)
    t0 = time.time()
    prs = Presentation(str(src))

    target = None  # (slide_idx, shape_idx, shape)
    slides = enumerate(prs.slides)
    for s_idx, slide in slides:
        if slide_index is not None and s_idx != slide_index:
            continue
        for sh_idx, shape in enumerate(slide.shapes):
            if not shape.has_text_frame:
                continue
            full_text = "\n".join(p.text for p in shape.text_frame.paragraphs)
            if find_text in full_text:
                target = (s_idx, sh_idx, shape)
                break
        if target:
            break

    if target is None:
        raise ValueError(f"no textbox found containing: {find_text!r}")

    s_idx, sh_idx, shape = target
    changes: dict[str, Any] = {}

    # Geometry
    if left_inches is not None:
        shape.left = Inches(left_inches); changes["left_inches"] = left_inches
    if top_inches is not None:
        shape.top = Inches(top_inches); changes["top_inches"] = top_inches
    if width_inches is not None:
        shape.width = Inches(width_inches); changes["width_inches"] = width_inches
    if height_inches is not None:
        shape.height = Inches(height_inches); changes["height_inches"] = height_inches

    # Text content replacement (preserves first run's style for new_text)
    if new_text is not None:
        tf = shape.text_frame
        # Clear existing paragraphs except first; then set first paragraph text
        # Use \n to split into paragraphs
        lines = new_text.split("\n")
        # Remove all but first paragraph
        first_p = tf.paragraphs[0]
        # Clear runs of first paragraph
        for r in list(first_p.runs):
            r._r.getparent().remove(r._r)
        # Remove subsequent paragraphs
        for p in list(tf.paragraphs[1:]):
            p._p.getparent().remove(p._p)
        # First line into first paragraph
        first_run = first_p.add_run()
        first_run.text = lines[0]
        # Additional lines as new paragraphs
        for line in lines[1:]:
            new_p = tf.add_paragraph()
            new_p.add_run().text = line
        changes["new_text"] = new_text[:80]

    # Font properties — apply to ALL runs across all paragraphs
    if font_size_pt is not None or bold is not None or italic is not None or font_color:
        tf = shape.text_frame
        for para in tf.paragraphs:
            for run in para.runs:
                if font_size_pt is not None:
                    run.font.size = Pt(font_size_pt)
                if bold is not None:
                    run.font.bold = bold
                if italic is not None:
                    run.font.italic = italic
                if font_color:
                    hex_c = font_color.lstrip("#").upper()
                    run.font.color.rgb = RGBColor(
                        int(hex_c[0:2], 16), int(hex_c[2:4], 16), int(hex_c[4:6], 16),
                    )
        if font_size_pt is not None:
            changes["font_size_pt"] = font_size_pt
        if bold is not None:
            changes["bold"] = bold
        if italic is not None:
            changes["italic"] = italic
        if font_color:
            changes["font_color"] = font_color

    fid, dst = _new_pptx_output()
    prs.save(str(dst))
    return {
        "file_id": fid, "url": public_url(fid),
        "modified_slide_index": s_idx,
        "modified_shape_index": sh_idx,
        "matched_text_preview": find_text[:60],
        "changes": changes,
        "ms_elapsed": int((time.time() - t0) * 1000),
    }


async def pptx_render_slide(
    pool, file_id: str, slide_index: Optional[int] = None, dpi: int = 110,
) -> dict[str, Any]:
    """Render PPTX slide(s) to PNG so the agent can VISUALLY verify the layout
    via describe_image. Pass slide_index=N for one slide, or omit to render all.
    Returns image_file_id(s) ready for describe_image."""
    import fitz
    from app.tools.convert import convert as _convert

    src = _resolve(file_id)
    t0 = time.time()

    # Step 1: PPTX → PDF (LO does the actual rendering)
    pdf_out = await _convert(pool, file_id=file_id, target_format="pdf")
    pdf_fid = pdf_out["file_id"]
    pdf_path = _resolve(pdf_fid)

    # Step 2: PyMuPDF rasterizes each PDF page → PNG
    doc = fitz.open(str(pdf_path))
    n = doc.page_count
    zoom = dpi / 72.0
    matrix = fitz.Matrix(zoom, zoom)

    def _render_one(idx: int) -> dict[str, Any]:
        if idx < 0 or idx >= n:
            raise ValueError(f"slide_index {idx} out of range (deck has {n} slides)")
        page = doc[idx]
        pix = page.get_pixmap(matrix=matrix)
        png_fid = new_file_id("png")
        with open(output_path(png_fid), "wb") as f:
            f.write(pix.tobytes("png"))
        return {
            "slide_index": idx,
            "image_file_id": png_fid,
            "url": public_url(png_fid),
            "width_px": pix.width,
            "height_px": pix.height,
        }

    if slide_index is not None:
        result = _render_one(slide_index)
        doc.close()
        result["ms_elapsed"] = int((time.time() - t0) * 1000)
        return result

    # Render all
    slides = [_render_one(i) for i in range(n)]
    doc.close()
    return {
        "slide_count": n,
        "slides": slides,
        "ms_elapsed": int((time.time() - t0) * 1000),
    }


async def pptx_slide_count(file_id: str) -> dict[str, Any]:
    """Quick count of slides (and their titles) in a deck."""
    from pptx import Presentation

    src = _resolve(file_id)
    prs = Presentation(str(src))
    titles = []
    for s in prs.slides:
        try:
            t = s.shapes.title.text_frame.text if s.shapes.title else ""
        except Exception:
            t = ""
        titles.append(t)
    return {"file_id": file_id, "slides": len(prs.slides), "titles": titles}


# Default catches: AI placeholder text + python-pptx layout boilerplate +
# common typos. Case-insensitive. Each entry is a regex string.
_DEFAULT_PLACEHOLDER_PATTERNS = [
    r"\bxxxx+\b",
    r"\blorem\s+ipsum\b",
    r"\bipsum\b",
    r"click\s+(?:here\s+)?to\s+(?:add|edit)",
    r"\btbd\b",
    r"\btodo\b",
    r"\bplaceholder\b",
    r"\bsample\s+text\b",
    r"your\s+(?:title|subtitle|text)\s+here",
    r"\[(?:organization|company|brand|date|name)[^\]]*\]",  # [organization name]
    r"<<[^>]+>>",                                            # <<replace_me>>
    r"\{\{[^}]+\}\}",                                        # {{handlebars}}
]


async def pptx_check_placeholders(
    file_id: str,
    extra_patterns: Optional[list[str]] = None,
) -> dict[str, Any]:
    """Scan all slide text for leftover placeholder markers (xxxx, lorem ipsum,
    "click to add", `[organization name]`, etc.). Use after building a deck —
    the visual self-review can miss text that LOOKS plausible but is template
    boilerplate the agent forgot to fill in. Returns per-slide hits with the
    matched substring and surrounding context."""
    import re as _re
    from pptx import Presentation

    src = _resolve(file_id)
    prs = Presentation(str(src))
    patterns = list(_DEFAULT_PLACEHOLDER_PATTERNS)
    if extra_patterns:
        patterns.extend(extra_patterns)
    compiled = [_re.compile(p, _re.IGNORECASE) for p in patterns]

    findings: list[dict[str, Any]] = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            try:
                text = shape.text_frame.text
            except Exception:
                continue
            if not text:
                continue
            for rx in compiled:
                m = rx.search(text)
                if not m:
                    continue
                start = max(0, m.start() - 20)
                end = min(len(text), m.end() + 20)
                findings.append({
                    "slide_index": slide_idx,
                    "shape_name": getattr(shape, "name", ""),
                    "matched": m.group(0),
                    "pattern": rx.pattern,
                    "context": text[start:end].replace("\n", " "),
                })
    return {
        "file_id": file_id,
        "slides_checked": len(prs.slides),
        "issues_found": len(findings),
        "findings": findings,
        "next_step_hint": (
            "Use pptx_modify_textbox(find_text=<matched>, new_text=<replacement>) "
            "to fix each finding, then call this tool again to verify clean."
            if findings else "All clean — no placeholder text detected."
        ),
    }
